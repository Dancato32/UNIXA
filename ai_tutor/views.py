from django.shortcuts import render, redirect, get_object_or_404
from django.urls import reverse
from django.contrib.auth.decorators import login_required
from django.contrib import messages
from django.http import JsonResponse, HttpResponse
from django.views.decorators.csrf import csrf_exempt
from django.core.paginator import Paginator
from django.utils import timezone
from django.views.decorators.http import require_POST
from .models import Conversation, EssayRequest, ChatThread
from .forms import ChatForm, EssayForm
from .ai_utils import ask_ai, deep_web_essay, generate_essay_with_sources, text_to_speech, resolve_model, AVAILABLE_MODELS, DEFAULT_MODEL
from nexa.rate_limit import ai_rate_limit, chat_rate_limit, essay_rate_limit, upload_rate_limit
import json
import os


def get_models(request):
    """Return available LLM models as JSON for the model picker UI."""
    return JsonResponse({'models': AVAILABLE_MODELS, 'default': DEFAULT_MODEL})


@login_required
def chat_ai(request, thread_id=None):
    """AI Chat interface - displays chat with past conversations and input form."""
    threads = ChatThread.objects.filter(user=request.user).order_by('-updated_at')
    
    current_thread = None
    conversations = []
    
    if thread_id:
        current_thread = get_object_or_404(ChatThread, id=thread_id, user=request.user)
        conversations = current_thread.messages.all().order_by('created_at')
    
    return render(request, 'ai_tutor/chat.html', {
        'threads': threads,
        'current_thread': current_thread,
        'conversations': conversations,
        'title': 'AI Tutor - CLUTCH'
    })

@login_required
@csrf_exempt
def create_thread(request):
    """Create a new chat thread and redirect to it."""
    thread = ChatThread.objects.create(user=request.user, title="New Chat")
    return JsonResponse({'success': True, 'thread_id': thread.id, 'url': f'/ai/chat/{thread.id}/'})

@login_required
@csrf_exempt
def delete_chat_thread(request, thread_id):
    """Delete a specific chat thread."""
    thread = get_object_or_404(ChatThread, id=thread_id, user=request.user)
    thread.delete()
    return JsonResponse({'success': True})

@login_required
@csrf_exempt
def rename_chat_thread(request, thread_id):
    """Rename a chat thread."""
    # - [x] Update `chat_stream` history limit to 50 messages (thread-specific)
    # - [x] Update `web_search_ajax` to use thread-specific history
    if request.method == 'POST':
        data = json.loads(request.body)
        title = data.get('title', 'New Chat').strip()[:100]
        thread = get_object_or_404(ChatThread, id=thread_id, user=request.user)
        thread.title = title
        thread.save()
        return JsonResponse({'success': True})
    return JsonResponse({'success': False}, status=405)


@login_required
def chat_ajax(request):
    """AJAX endpoint for real-time chat without page reload."""
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            message = data.get('message', '').strip()
            use_rag = data.get('use_rag', True)
            model = resolve_model(data.get('model'))
            
            if not message:
                return JsonResponse({'error': 'Message cannot be empty'}, status=400)
            
            thread_id = data.get('thread_id')
            if thread_id:
                thread = get_object_or_404(ChatThread, id=thread_id, user=request.user)
            else:
                thread = ChatThread.objects.create(user=request.user, title=message[:50])

            # Fetch optimized history for session context (Infinity Memory - Balanced)
            recent = Conversation.objects.filter(thread=thread).order_by('-created_at')[:20]
            history = [{'message': c.message, 'response': c.response} for c in reversed(recent)]

            # Get AI response with history
            response = ask_ai(message, user=request.user, use_rag=use_rag, history=history, model=model)
            
            conversation = Conversation.objects.create(
                user=request.user,
                thread=thread,
                message=message,
                response=response
            )
            
            # Update thread title if it's the first message
            if thread.messages.count() == 1:
                thread.title = message[:50]
                thread.save()
            
            return JsonResponse({
                'success': True,
                'message': message,
                'response': response,
                'timestamp': conversation.created_at.strftime('%H:%M')
            })
            
        except json.JSONDecodeError:
            return JsonResponse({'error': 'Invalid JSON'}, status=400)
        except Exception as e:
            return JsonResponse({'error': str(e)}, status=500)
    
    return JsonResponse({'error': 'Method not allowed'}, status=405)


@login_required
def text_to_speech_view(request):
    """Convert text to speech using Resemble.ai, return base64 data URL â€” works on all platforms including mobile."""
    if request.method == 'POST':
        try:
            import base64 as _b64
            data = json.loads(request.body)
            text = data.get('text', '')

            if not text:
                return JsonResponse({'error': 'Text cannot be empty'}, status=400)

            audio_content = text_to_speech(text)

            if audio_content:
                # Return as base64 data URL â€” no file system, no media serving needed
                # Works on iOS Safari, Android Chrome, all mobile browsers
                b64 = _b64.b64encode(audio_content).decode('utf-8')
                audio_url = f'data:audio/mpeg;base64,{b64}'
                return JsonResponse({'audio_url': audio_url})
            else:
                return JsonResponse({'use_browser_tts': True})

        except Exception as e:
            return JsonResponse({'error': str(e)}, status=500)

    return JsonResponse({'error': 'Method not allowed'}, status=405)


@login_required
def essay_request(request):
    """Essay request interface - form for topic and list of past essays."""
    if request.method == 'POST':
        form = EssayForm(request.POST)
        if form.is_valid():
            topic = form.cleaned_data['topic']
            output_format = form.cleaned_data.get('output_format', 'text')
            word_count = int(form.cleaned_data.get('word_count') or 500)
            writing_style = request.POST.get('writing_style', 'argumentative')

            # Build topic string with style hint for the generator
            styled_topic = f"{topic} [Style: {writing_style}]"

            # Generate essay with research
            essay_text, sources = generate_essay_with_sources(styled_topic, user=request.user, word_count=word_count)

            # Save essay request
            essay = EssayRequest.objects.create(
                user=request.user,
                topic=topic,
                word_count=word_count,
                research_done=True,
                essay_text=essay_text,
                output_format=output_format
            )

            messages.success(request, f'Essay generated on "{topic}"!')
            return redirect('essay_detail', essay_id=essay.id)
    else:
        form = EssayForm()
    
    # Get user's past essays
    essays = EssayRequest.objects.filter(user=request.user)
    paginator = Paginator(essays, 10)
    page_number = request.GET.get('page')
    essays_page = paginator.get_page(page_number)
    
    return render(request, 'ai_tutor/essay_request.html', {
        'form': form,
        'essays': essays_page,
        'title': 'AI Tutor - Essay Generator'
    })


@login_required
@require_POST
def essay_generate_ajax(request):
    """
    Elite Studio AJAX Generation.
    Handles the high-fidelity multi-step generation process.
    """
    try:
        data = json.loads(request.body)
        topic = data.get('topic', '').strip()
        word_count = int(data.get('word_count') or 500)
        writing_style = data.get('style', 'argumentative')
        use_research = data.get('research', True)
        
        if not topic:
            return JsonResponse({'ok': False, 'error': 'Topic is required.'})

        # Create a blank database record to hold the streamed text later
        essay = EssayRequest.objects.create(
            user=request.user,
            topic=topic,
            word_count=word_count,
            research_done=use_research,
            essay_text="",  # Empty, will be filled by stream
            output_format='text'
        )
        
        # Store metadata intended for stream generation in session
        # so essay_stream_build can use it
        rich_topic = f"{topic} [Style: {writing_style}] [Professional Structure: Yes] [Humanized: Yes]"
        request.session[f'essay_build_topic_{essay.id}'] = rich_topic
        request.session[f'essay_build_wc_{essay.id}'] = word_count
        request.session[f'essay_build_research_{essay.id}'] = use_research
        
        redirect_url = reverse('essay_detail', args=[essay.id]) + '?stream=true'

        return JsonResponse({
            'ok': True, 
            'essay_id': essay.id,
            'redirect_url': redirect_url
        })
    except Exception as e:
        return JsonResponse({'ok': False, 'error': str(e)})


@login_required
def essay_stream_build(request, essay_id):
    """
    Stream the essay generation directly to the client using SSE.
    """
    from django.http import StreamingHttpResponse
    from .ai_utils import generate_essay_stream
    
    essay = get_object_or_404(EssayRequest, id=essay_id, user=request.user)
    
    # Retrieve generation parameters from session
    topic = request.session.get(f'essay_build_topic_{essay.id}', essay.topic)
    word_count = request.session.get(f'essay_build_wc_{essay.id}', essay.word_count)
    use_research = request.session.get(f'essay_build_research_{essay.id}', essay.research_done)
    model = request.GET.get('model', None)
    
    def event_stream():
        full_text = []
        try:
            for chunk in generate_essay_stream(topic, request.user, word_count, model=model, use_research=use_research):
                full_text.append(chunk)
                # SSE format
                # Replace newlines with a unique token or just use JSON
                import json
                yield f"data: {json.dumps({'chunk': chunk})}\n\n"
                
            # Once stream completes, save to database
            final_text = "".join(full_text)
            essay.essay_text = final_text
            essay.save()
            
            # Optionally clean up session keys
            try:
                del request.session[f'essay_build_topic_{essay.id}']
                del request.session[f'essay_build_wc_{essay.id}']
            except:
                pass
                
            yield "event: done\ndata: {}\n\n"
        except Exception as e:
            yield f"event: error\ndata: {json.dumps({'error': str(e)})}\n\n"
            
    return StreamingHttpResponse(event_stream(), content_type='text/event-stream')



@login_required
def essay_edit_chat(request):
    """AJAX: chat with AI to edit/improve the current essay. Returns a reply and optionally an updated essay."""
    if request.method != 'POST':
        return JsonResponse({'ok': False, 'error': 'POST required'}, status=405)
    try:
        data = json.loads(request.body)
        message = data.get('message', '').strip()
        essay_text = data.get('essay_text', '').strip()
        if not message or not essay_text:
            return JsonResponse({'ok': False, 'error': 'Message and essay required'}, status=400)

        model = resolve_model(data.get('model'))

        from .ai_utils import get_openai_client, build_system_prompt
        client = get_openai_client()

        system = (
            "You are Nexa, an expert essay editor. The user will share their essay and ask you to edit or improve it. "
            "If the user asks for a change (e.g. 'make it more persuasive', 'shorten it', 'fix grammar'), "
            "return BOTH a short conversational reply AND the full updated essay. "
            "If the user asks a question or wants advice without a rewrite, just reply conversationally. "
            "Format your response as JSON: "
            '{"reply": "short message to user", "updated_essay": "full updated essay or null if no rewrite needed"} '
            "Return ONLY valid JSON, no extra text."
        )

        model_id = resolve_model(data.get('model'))
        resp = client.chat.completions.create(
            model=model_id,
            messages=[
                {"role": "system", "content": system},
                {"role": "user", "content": f"My essay:\n\n{essay_text}\n\nMy request: {message}"},
            ],
            max_tokens=2000,
            temperature=0.7,
        )
        raw = resp.choices[0].message.content.strip()

        # Parse JSON response
        import re
        match = re.search(r'\{.*\}', raw, re.DOTALL)
        if match:
            result = json.loads(match.group())
        else:
            result = {'reply': raw, 'updated_essay': None}

        return JsonResponse({
            'ok': True,
            'reply': result.get('reply', ''),
            'updated_essay': result.get('updated_essay') or None,
        })
    except Exception as e:
        return JsonResponse({'ok': False, 'error': str(e)}, status=500)


@login_required
def essay_generate_ajax(request):
    """AJAX: generate essay and return JSON so the frontend can stream/typewrite it."""
    if request.method != 'POST':
        return JsonResponse({'ok': False, 'error': 'POST required'}, status=405)
    try:
        data = json.loads(request.body)
        topic = data.get('topic', '').strip()
        if not topic:
            return JsonResponse({'ok': False, 'error': 'Topic is required'}, status=400)
        word_count = int(data.get('word_count', 500))
        writing_style = data.get('writing_style', 'argumentative')
        output_format = data.get('output_format', 'text')

        styled_topic = f"{topic} [Style: {writing_style}]"

        # Inject selected references into the prompt
        references = data.get('references', [])  # list of {title, snippet, url}
        ref_block = ''
        if references:
            ref_lines = '\n'.join(
                f"- [{r['title']}]({r['url']}): {r['snippet']}" for r in references
            )
            ref_block = f"\n\nUse the following real sources as references and cite them naturally in the essay:\n{ref_lines}"
            styled_topic += ref_block

        essay_text, sources = generate_essay_with_sources(styled_topic, user=request.user, word_count=word_count)

        essay = EssayRequest.objects.create(
            user=request.user,
            topic=topic,
            word_count=word_count,
            research_done=True,
            essay_text=essay_text,
            output_format=output_format,
        )
        return JsonResponse({
            'ok': True,
            'essay_text': essay_text,
            'essay_id': essay.id,
            'topic': topic,
            'detail_url': f'/ai/essay/{essay.id}/',
        })
    except Exception as e:
        return JsonResponse({'ok': False, 'error': str(e)}, status=500)


@login_required
def essay_detail(request, essay_id):

    """View essay details."""
    essay = get_object_or_404(EssayRequest, id=essay_id, user=request.user)
    
    return render(request, 'ai_tutor/essay_detail.html', {
        'essay': essay,
        'title': f'Essay: {essay.topic}'
    })


@login_required
def delete_essay(request, essay_id):
    """Delete an essay request."""
    try:
        essay = EssayRequest.objects.get(id=essay_id, user=request.user)
    except EssayRequest.DoesNotExist:
        messages.error(request, 'Essay not found or already deleted.')
        return redirect('essay_request')
    
    if request.method == 'POST':
        essay.delete()
        messages.success(request, 'Essay deleted successfully.')
        return redirect('essay_request')
    
    return render(request, 'ai_tutor/essay_delete.html', {
        'essay': essay,
        'title': 'Delete Essay'
    })


@login_required
def chat_with_image(request):
    """SSE streaming endpoint for Multi-Modal vision and problem solving."""
    if request.method != 'POST':
        return JsonResponse({'error': 'Method not allowed'}, status=405)

    try:
        message = request.POST.get('message', '').strip() or 'Explain this image step by step.'
        image_file = request.FILES.get('image')
        use_rag = request.POST.get('use_rag', 'true').lower() == 'true'

        if not image_file:
            return JsonResponse({'error': 'Please select an image first.'}, status=400)
            
        if image_file.size > 4 * 1024 * 1024:
            return JsonResponse({'error': 'Image is too large (max 4MB).'}, status=400)

        import base64
        import traceback
        from .ai_utils import get_openai_client, build_system_prompt
        from django.http import StreamingHttpResponse

        # Resolve History (Session Context)
        thread_id = request.POST.get('thread_id')
        if thread_id:
            thread = get_object_or_404(ChatThread, id=thread_id, user=request.user)
        else:
            thread = ChatThread.objects.create(user=request.user, title=message[:50])

        recent = Conversation.objects.filter(thread=thread).order_by('-created_at')[:100]
        # Evaluate QuerySet to avoid any late evaluation issues in the generator
        history_list = list(reversed(recent))
        history = [{'message': c.message, 'response': c.response} for c in history_list]

        # Reset pointer and read
        image_file.seek(0)
        image_data = base64.b64encode(image_file.read()).decode('utf-8')
        mime = image_file.content_type or 'image/jpeg'

        client = get_openai_client()
        system_message = build_system_prompt(use_rag=use_rag, user=request.user, is_vision=True, query=message)
        
        messages_list = [{"role": "system", "content": system_message}]
        for h in history:
            text_context = str(h['message'] or '').replace('[Image Analysis] ', '').replace('[Web Search] ', '')
            messages_list.append({"role": "user", "content": text_context})
            messages_list.append({"role": "assistant", "content": str(h['response'] or '')})
            
        messages_list.append({
            "role": "user",
            "content": [
                {"type": "text", "text": message},
                {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{image_data}"}}
            ]
        })

        def event_stream():
            full_response = []
            try:
                # Switching back to gpt-4o-mini as it is the standard vision model on OpenRouter
                stream = client.chat.completions.create(
                    model="openai/gpt-4o-mini",
                    messages=messages_list,
                    max_tokens=2000,
                    stream=True,
                )
                for chunk in stream:
                    if not chunk.choices: continue
                    token = chunk.choices[0].delta.content or ''
                    if token:
                        full_response.append(token)
                        yield f"data: {json.dumps({'token': token})}\n\n"

                complete = ''.join(full_response)
                if complete:
                    Conversation.objects.create(
                        user=request.user,
                        message=f"[Image Analysis] {message}",
                        response=complete
                    )
                yield f"data: {json.dumps({'done': True})}\n\n"
            except Exception as stream_err:
                yield f"data: {json.dumps({'error': str(stream_err)})}\n\n"

        response = StreamingHttpResponse(event_stream(), content_type='text/event-stream')
        response['Cache-Control'] = 'no-cache, no-transform'
        response['X-Accel-Buffering'] = 'no'
        return response

    except Exception as e:
        import traceback
        print(f"Vision Error: {str(e)}\n{traceback.format_exc()}")
        return JsonResponse({'error': f"Processing failed: {str(e)}"}, status=500)


@login_required
def essay_web_search(request):
    """Search the web and return structured article cards for the user to pick as references."""
    if request.method != 'POST':
        return JsonResponse({'error': 'POST required'}, status=405)
    try:
        data = json.loads(request.body)
        query = data.get('query', '').strip()
        if not query:
            return JsonResponse({'error': 'Query required'}, status=400)

        import urllib.request, urllib.parse, re

        results = []

        # â”€â”€ DuckDuckGo Instant Answer API â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
        ddg_url = f"https://api.duckduckgo.com/?q={urllib.parse.quote(query)}&format=json&no_html=1&skip_disambig=1"
        req = urllib.request.Request(ddg_url, headers={'User-Agent': 'Mozilla/5.0'})
        with urllib.request.urlopen(req, timeout=6) as resp:
            ddg = json.loads(resp.read().decode())

        # Abstract (Wikipedia-style summary)
        if ddg.get('AbstractText') and ddg.get('AbstractURL'):
            results.append({
                'title': ddg.get('Heading') or query,
                'snippet': ddg['AbstractText'][:300],
                'url': ddg['AbstractURL'],
                'source': ddg.get('AbstractSource', 'Web'),
            })

        # Related topics
        for r in ddg.get('RelatedTopics', [])[:8]:
            if isinstance(r, dict) and r.get('Text') and r.get('FirstURL'):
                title = r.get('Text', '')[:60]
                results.append({
                    'title': title,
                    'snippet': r['Text'][:300],
                    'url': r['FirstURL'],
                    'source': 'DuckDuckGo',
                })
            # Handle grouped topics
            elif isinstance(r, dict) and r.get('Topics'):
                for sub in r['Topics'][:3]:
                    if sub.get('Text') and sub.get('FirstURL'):
                        results.append({
                            'title': sub['Text'][:60],
                            'snippet': sub['Text'][:300],
                            'url': sub['FirstURL'],
                            'source': 'DuckDuckGo',
                        })

        # Results section (news/web results)
        for r in ddg.get('Results', [])[:5]:
            if r.get('Text') and r.get('FirstURL'):
                results.append({
                    'title': r.get('Text', '')[:80],
                    'snippet': r['Text'][:300],
                    'url': r['FirstURL'],
                    'source': 'Web',
                })

        # Deduplicate by URL
        seen = set()
        unique = []
        for r in results:
            if r['url'] not in seen:
                seen.add(r['url'])
                unique.append(r)

        return JsonResponse({'ok': True, 'results': unique[:10], 'query': query})
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


@login_required
def web_search_ajax(request):
    """AJAX endpoint: perform a web search then answer with context."""
    if request.method != 'POST':
        return JsonResponse({'error': 'Method not allowed'}, status=405)

    try:
        data = json.loads(request.body)
        query = data.get('message', '').strip()

        if not query:
            return JsonResponse({'error': 'Query cannot be empty'}, status=400)

        # Fetch search results via DuckDuckGo (no API key needed)
        import urllib.request
        import urllib.parse
        import re

        search_url = f"https://api.duckduckgo.com/?q={urllib.parse.quote(query)}&format=json&no_html=1&skip_disambig=1"
        req = urllib.request.Request(search_url, headers={'User-Agent': 'Mozilla/5.0'})
        with urllib.request.urlopen(req, timeout=5) as resp:
            ddg = json.loads(resp.read().decode())

        snippets = []
        if ddg.get('AbstractText'):
            snippets.append(ddg['AbstractText'])
        for r in ddg.get('RelatedTopics', [])[:4]:
            if isinstance(r, dict) and r.get('Text'):
                snippets.append(r['Text'])

        web_context = '\n'.join(snippets[:5]) if snippets else ''

        client = get_openai_client()

        system = """You are Nexa, an AI tutor. The user asked a question and web search results are provided below.
Use the search results to give an accurate, up-to-date answer. Cite facts from the results naturally.
Maintain consistency with the previous conversation in this thread.
Format your answer in clear steps or paragraphs as appropriate. Use LaTeX for any math."""

        thread_id = data.get('thread_id')
        if thread_id:
            thread = get_object_or_404(ChatThread, id=thread_id, user=request.user)
            recent = Conversation.objects.filter(thread=thread).order_by('-created_at')[:50]
        else:
            thread = None
            recent = Conversation.objects.filter(user=request.user, thread__isnull=True).order_by('-created_at')[:50]
            
        history = [{'message': c.message, 'response': c.response} for c in reversed(recent)]

        user_content = f"Question: {query}"
        if web_context:
            user_content += f"\n\nWeb search results:\n{web_context}"

        messages_list = [{"role": "system", "content": system}]
        for h in history:
            messages_list.append({"role": "user", "content": h['message']})
            messages_list.append({"role": "assistant", "content": h['response']})
        messages_list.append({"role": "user", "content": user_content})

        model_ws = resolve_model(data.get('model'))

        response = client.chat.completions.create(
            model=model_ws,
            messages=messages_list,
            max_tokens=1000,
            temperature=0.7
        )
        ai_response = response.choices[0].message.content

        conversation = Conversation.objects.create(
            user=request.user,
            thread=thread,
            message=f"[Web Search] {query}",
            response=ai_response
        )

        return JsonResponse({
            'success': True,
            'response': ai_response,
            'timestamp': conversation.created_at.strftime('%H:%M'),
            'web_used': bool(web_context)
        })

    except json.JSONDecodeError:
        return JsonResponse({'error': 'Invalid JSON'}, status=400)
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


@login_required
def chat_stream(request):
    """SSE streaming endpoint â€” streams AI response token by token."""
    if request.method != 'POST':
        return JsonResponse({'error': 'POST required'}, status=405)

    try:
        data = json.loads(request.body)
        message = data.get('message', '').strip()
        use_rag = data.get('use_rag', True)

        model = resolve_model(data.get('model'))

        if not message:
            return JsonResponse({'error': 'Empty message'}, status=400)

        from .ai_utils import get_openai_client, get_study_materials_for_rag, build_system_prompt
        from django.http import StreamingHttpResponse

        thread_id = data.get('thread_id')
        if thread_id:
            thread = get_object_or_404(ChatThread, id=thread_id, user=request.user)
        else:
            thread = ChatThread.objects.create(user=request.user, title=message[:50])

        # Fetch optimized history for session context (Infinity Memory - Balanced)
        recent = Conversation.objects.filter(thread=thread).order_by('-created_at')[:50]
        history = [{'message': c.message, 'response': c.response} for c in reversed(recent)]

        client = get_openai_client()
        system_message = build_system_prompt(use_rag, request.user, query=message)

        messages_list = [{"role": "system", "content": system_message}]
        for h in history:
            messages_list.append({"role": "user", "content": h['message']})
            messages_list.append({"role": "assistant", "content": h['response']})
        messages_list.append({"role": "user", "content": message})

        def event_stream():
            full_response = []
            try:
                stream = client.chat.completions.create(
                    model=model,
                    messages=messages_list,
                    max_tokens=1000,
                    temperature=0.7,
                    stream=True,
                )
                for chunk in stream:
                    token = chunk.choices[0].delta.content or ''
                    if token:
                        full_response.append(token)
                        # SSE format: data: <token>\n\n
                        yield f"data: {json.dumps({'token': token})}\n\n"

                # Save full response to DB
                complete = "".join(full_response)
                Conversation.objects.create(
                    user=request.user,
                    thread=thread,
                    message=message,
                    response=complete,
                )
                
                # Update title if it's the first message
                if thread.messages.count() == 1:
                    thread.title = message[:50]
                    thread.save()
                yield f"data: {json.dumps({'done': True, 'thread_id': thread.id})}\n\n"

            except Exception as e:
                yield f"data: {json.dumps({'error': str(e)})}\n\n"

        response = StreamingHttpResponse(event_stream(), content_type='text/event-stream')
        response['Cache-Control'] = 'no-cache, no-transform'
        response['X-Accel-Buffering'] = 'no'
        return response

    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


@login_required
def clear_conversations(request):
    """Clear all chat conversations for the user."""
    if request.method == 'POST':
        Conversation.objects.filter(user=request.user).delete()
        messages.success(request, 'All conversations cleared.')
        return redirect('ai_chat')
    return redirect('ai_chat')


@login_required
def essay_guidance(request):
    """Smart Essay Guidance â€” topic breakdown, thesis, outline before writing."""
    if request.method != 'POST':
        return JsonResponse({'error': 'Method not allowed'}, status=405)
    try:
        data = json.loads(request.body)
        topic = data.get('topic', '').strip()
        action = data.get('action', 'all')  # breakdown | thesis | outline | all
        if not topic:
            return JsonResponse({'error': 'Topic required'}, status=400)

        # Build reference context block
        references = data.get('references', [])
        ref_block = ''
        if references:
            ref_lines = '\n'.join(f"- [{r['title']}]({r['url']}): {r['snippet']}" for r in references)
            ref_block = f"\n\nReal-world sources to draw from:\n{ref_lines}\n"

        from .ai_utils import get_openai_client
        client = get_openai_client()

        prompts = {
            'breakdown': (
                f"A student wants to write an essay on: \"{topic}\"\n\n"
                "Give a clear topic breakdown:\n"
                "1. What the question is really asking (in plain language)\n"
                "2. Key concepts and terms they must understand\n"
                "3. Common mistakes students make on this topic\n"
                "4. What a strong essay on this topic must include\n\n"
                f"{ref_block}"
                "Be concise, practical, and student-friendly. No markdown."
            ),
            'thesis': (
                f"Topic: \"{topic}\"\n\n"
                "Generate exactly 3 strong thesis statements for this essay topic.\n"
                "Each thesis should:\n"
                "- Make a clear, arguable claim\n"
                "- Be one sentence\n"
                "- Be suitable for an academic essay\n\n"
                f"{ref_block}"
                "Format:\nThesis 1: ...\nThesis 2: ...\nThesis 3: ...\n\nNo extra commentary."
            ),
            'outline': (
                f"Topic: \"{topic}\"\n\n"
                "Build a detailed essay outline with:\n"
                "Introduction: Hook idea, background context, thesis placeholder\n"
                "Body Paragraph 1: Main point, supporting evidence, explanation\n"
                "Body Paragraph 2: Main point, supporting evidence, explanation\n"
                "Body Paragraph 3: Main point, supporting evidence, explanation\n"
                "Conclusion: Restate thesis, summarise points, closing thought\n\n"
                f"{ref_block}"
                "Be specific to this topic. No markdown symbols."
            ),
        }

        model_m = resolve_model(data.get('model'))

        if action == 'all':
            # Run all three and combine
            results = {}
            for key, prompt in prompts.items():
                resp = client.chat.completions.create(
                    model=model_m,
                    messages=[
                        {"role": "system", "content": "You are an expert academic writing coach helping students plan essays. Be clear, practical, and encouraging."},
                        {"role": "user", "content": prompt}
                    ],
                    max_tokens=600, temperature=0.6,
                )
                results[key] = resp.choices[0].message.content.strip()
            return JsonResponse({'ok': True, 'results': results})
        else:
            prompt = prompts.get(action, prompts['breakdown'])
            resp = client.chat.completions.create(
                model=model_m,
                messages=[
                    {"role": "system", "content": "You are an expert academic writing coach. Be clear, practical, and encouraging."},
                    {"role": "user", "content": prompt}
                ],
                max_tokens=600, temperature=0.6,
            )
            return JsonResponse({'ok': True, 'result': resp.choices[0].message.content.strip(), 'action': action})
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


@login_required
def essay_improve(request):
    """Essay Improver â€” grammar, vocabulary, clarity, tone adjustment."""
    if request.method != 'POST':
        return JsonResponse({'error': 'Method not allowed'}, status=405)
    try:
        data = json.loads(request.body)
        essay_text = data.get('essay_text', '').strip()
        improvements = data.get('improvements', ['grammar', 'clarity'])
        tone = data.get('tone', '')
        references = data.get('references', [])
        if not essay_text:
            return JsonResponse({'error': 'Essay text required'}, status=400)

        improvement_instructions = []
        if 'grammar' in improvements:
            improvement_instructions.append("Fix all grammar, punctuation, and spelling errors")
        if 'vocabulary' in improvements:
            improvement_instructions.append("Upgrade vocabulary with more precise and varied word choices")
        if 'clarity' in improvements:
            improvement_instructions.append("Improve sentence clarity and flow â€” break up run-ons, vary sentence length")
        if 'structure' in improvements:
            improvement_instructions.append("Improve paragraph structure and logical flow between ideas")
        if tone:
            improvement_instructions.append(f"Adjust the tone to be {tone}")
        if references:
            ref_lines = '\n'.join(f"- [{r['title']}]({r['url']}): {r['snippet']}" for r in references)
            improvement_instructions.append(f"Incorporate and cite these real sources where relevant:\n{ref_lines}")

        instructions_text = "\n".join(f"- {i}" for i in improvement_instructions)

        model_i = resolve_model(data.get('model'))

        from .ai_utils import get_openai_client
        client = get_openai_client()
        resp = client.chat.completions.create(
            model=model_i,
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You are an expert essay editor and writing coach. "
                        "Improve the student's essay according to the instructions. "
                        "Preserve their voice and all their ideas. "
                        "Output ONLY the improved essay â€” no commentary, no labels, no markdown."
                    )
                },
                {
                    "role": "user",
                    "content": f"Improvement instructions:\n{instructions_text}\n\nStudent's essay:\n{essay_text}"
                }
            ],
            max_tokens=2000, temperature=0.6,
        )
        improved = resp.choices[0].message.content.strip()

        # Also generate a brief feedback summary
        feedback_resp = client.chat.completions.create(
            model=model_i,
            messages=[
                {"role": "system", "content": "You are a writing coach. Give brief, encouraging feedback."},
                {
                    "role": "user",
                    "content": (
                        f"Original essay:\n{essay_text[:1000]}\n\n"
                        f"Improved essay:\n{improved[:1000]}\n\n"
                        "In 2-3 sentences, summarise the key improvements made. Be specific and encouraging."
                    )
                }
            ],
            max_tokens=150, temperature=0.5,
        )
        feedback = feedback_resp.choices[0].message.content.strip()
        return JsonResponse({'ok': True, 'improved': improved, 'feedback': feedback})
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


@login_required
def essay_restyle(request, essay_id):
    """AI restyle: rewrite essay text based on user style instructions."""
    if request.method != 'POST':
        return JsonResponse({'error': 'Method not allowed'}, status=405)
    essay = get_object_or_404(EssayRequest, id=essay_id, user=request.user)
    try:
        data = json.loads(request.body)
        instructions = data.get('instructions', '').strip()
        current_text = data.get('current_text', essay.essay_text).strip()
        if not instructions:
            return JsonResponse({'error': 'No instructions provided'}, status=400)

        from .ai_utils import get_openai_client
        client = get_openai_client()
        response = client.chat.completions.create(
            model=resolve_model(data.get('model')),
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You are an expert essay editor. The user will give you an essay and styling instructions. "
                        "Rewrite the essay following those instructions exactly. "
                        "Output ONLY the rewritten essay â€” no commentary, no headings, no markdown. "
                        "Preserve all facts and arguments. Only change style, tone, structure, or format as instructed."
                    )
                },
                {
                    "role": "user",
                    "content": f"Style instructions: {instructions}\n\nEssay to restyle:\n{current_text}"
                }
            ],
            max_tokens=2000,
            temperature=0.75,
        )
        new_text = response.choices[0].message.content.strip()
        return JsonResponse({'ok': True, 'text': new_text})
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


@login_required
def essay_save_edits(request, essay_id):
    """Save user edits to essay text."""
    if request.method != 'POST':
        return JsonResponse({'error': 'Method not allowed'}, status=405)
    essay = get_object_or_404(EssayRequest, id=essay_id, user=request.user)
    try:
        data = json.loads(request.body)
        content = data.get('content', '').strip()
        if content:
            essay.essay_text = content
            essay.save()
        return JsonResponse({'ok': True})
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


@login_required
def export_essay(request):
    """Export essay in the specified format."""
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            essay_text = data.get('essay_text', '')
            output_format = data.get('format', 'text')
            topic = data.get('topic', 'essay')
            
            if not essay_text:
                return JsonResponse({'error': 'No essay text provided'}, status=400)
            
            from io import BytesIO
            
            if output_format == 'word':
                from docx import Document
                doc = Document()
                doc.add_heading(topic, 0)
                for paragraph in essay_text.split('\n\n'):
                    if paragraph.strip():
                        doc.add_paragraph(paragraph.strip())
                
                buffer = BytesIO()
                doc.save(buffer)
                buffer.seek(0)
                response = HttpResponse(
                    buffer.getvalue(),
                    content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                )
                response['Content-Disposition'] = f'attachment; filename="{topic[:30]}.docx"'
                return response
            
            elif output_format == 'pdf':
                from reportlab.lib.pagesizes import letter
                from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
                from reportlab.lib.units import inch
                from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
                
                buffer = BytesIO()
                doc = SimpleDocTemplate(buffer, pagesize=letter, rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=18)
                styles = getSampleStyleSheet()
                title_style = ParagraphStyle('Title', parent=styles['Heading1'], fontSize=18, spaceAfter=12)
                body_style = ParagraphStyle('Body', parent=styles['Normal'], fontSize=11, spaceAfter=12)
                
                story = []
                story.append(Paragraph(topic, title_style))
                story.append(Spacer(1, 0.2*inch))
                
                for paragraph in essay_text.split('\n\n'):
                    if paragraph.strip():
                        story.append(Paragraph(paragraph.strip(), body_style))
                        story.append(Spacer(1, 0.1*inch))
                
                doc.build(story)
                buffer.seek(0)
                
                response = HttpResponse(buffer.getvalue(), content_type='application/pdf')
                response['Content-Disposition'] = f'attachment; filename="{topic[:30]}.pdf"'
                return response
            
            elif output_format == 'powerpoint':
                from pptx import Presentation
                from pptx.util import Inches, Pt, Emu
                from pptx.dml.color import RGBColor
                from pptx.enum.text import PP_ALIGN
                import urllib.request
                import urllib.parse
                import re as _re

                # â”€â”€ AI: generate structured slide outline â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
                from .ai_utils import get_openai_client
                ai_client = get_openai_client()
                outline_resp = ai_client.chat.completions.create(
                    model="openai/gpt-4o-mini",
                    messages=[
                        {
                            "role": "system",
                            "content": (
                                "You are a professional presentation designer. "
                                "Given an essay, produce a JSON array of slides. "
                                "Each slide: {\"title\": str, \"bullets\": [str, ...], \"image_query\": str}. "
                                "image_query is a short 2-4 word Unsplash search term for a relevant photo. "
                                "First slide is the title slide with no bullets. "
                                "Produce 6-8 slides total. Output ONLY valid JSON, no markdown."
                            )
                        },
                        {"role": "user", "content": f"Topic: {topic}\n\nEssay:\n{essay_text[:3000]}"}
                    ],
                    max_tokens=1200,
                    temperature=0.5,
                )
                raw_json = outline_resp.choices[0].message.content.strip()
                # strip markdown code fences if present
                raw_json = _re.sub(r'^```[a-z]*\n?', '', raw_json)
                raw_json = _re.sub(r'\n?```$', '', raw_json)
                try:
                    slides_data = json.loads(raw_json)
                except Exception:
                    # fallback: split essay into sections
                    sections = [s.strip() for s in essay_text.split('\n\n') if s.strip()]
                    slides_data = [{"title": topic, "bullets": [], "image_query": topic}]
                    for sec in sections[:7]:
                        lines = sec.split('\n')
                        slides_data.append({
                            "title": lines[0][:60],
                            "bullets": [l.strip() for l in lines[1:4] if l.strip()],
                            "image_query": topic
                        })

                # â”€â”€ Helper: fetch image bytes from Unsplash source â”€â”€â”€â”€â”€â”€â”€â”€â”€
                def fetch_image(query):
                    try:
                        safe = urllib.parse.quote(query)
                        url = f"https://source.unsplash.com/1280x720/?{safe}"
                        req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
                        with urllib.request.urlopen(req, timeout=8) as r:
                            return r.read()
                    except Exception:
                        return None

                # â”€â”€ Design constants â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
                DARK_BG   = RGBColor(0x0D, 0x0D, 0x0D)
                ACCENT    = RGBColor(0x6C, 0x63, 0xFF)   # purple
                WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
                LIGHT_TXT = RGBColor(0xCC, 0xCC, 0xCC)
                W = Inches(13.33)
                H = Inches(7.5)

                prs = Presentation()
                prs.slide_width  = W
                prs.slide_height = H

                def add_bg(slide, color=DARK_BG):
                    """Fill slide background with solid color."""
                    fill = slide.background.fill
                    fill.solid()
                    fill.fore_color.rgb = color

                def add_textbox(slide, text, left, top, width, height,
                                font_size=24, bold=False, color=WHITE,
                                align=PP_ALIGN.LEFT, wrap=True):
                    txb = slide.shapes.add_textbox(left, top, width, height)
                    tf  = txb.text_frame
                    tf.word_wrap = wrap
                    p   = tf.paragraphs[0]
                    p.alignment = align
                    run = p.add_run()
                    run.text = text
                    run.font.size = Pt(font_size)
                    run.font.bold = bold
                    run.font.color.rgb = color
                    return txb

                def add_accent_bar(slide, top=Inches(0.55)):
                    """Thin horizontal accent line."""
                    bar = slide.shapes.add_shape(
                        1,  # MSO_SHAPE_TYPE.RECTANGLE
                        Inches(0.5), top, Inches(12.33), Inches(0.04)
                    )
                    bar.fill.solid()
                    bar.fill.fore_color.rgb = ACCENT
                    bar.line.fill.background()

                # â”€â”€ Slide 1: Title slide â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
                title_data = slides_data[0] if slides_data else {"title": topic, "image_query": topic}
                img_bytes = fetch_image(title_data.get("image_query", topic))

                blank_layout = prs.slide_layouts[6]  # completely blank
                slide = prs.slides.add_slide(blank_layout)
                add_bg(slide)

                # Background image (dimmed)
                if img_bytes:
                    from io import BytesIO as _BIO
                    pic = slide.shapes.add_picture(_BIO(img_bytes), Inches(0), Inches(0), W, H)
                    # Send image to back
                    slide.shapes._spTree.remove(pic._element)
                    slide.shapes._spTree.insert(2, pic._element)
                    # Dark overlay
                    overlay = slide.shapes.add_shape(1, Inches(0), Inches(0), W, H)
                    overlay.fill.solid()
                    overlay.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
                    overlay.line.fill.background()
                    overlay.fill.fore_color.theme_color = None
                    # set transparency via XML
                    from lxml import etree
                    solidFill = overlay.fill._xPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
                    if solidFill is not None:
                        srgb = solidFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                        if srgb is not None:
                            alpha = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                            alpha.set('val', '75000')  # 75% opacity overlay

                add_accent_bar(slide, top=Inches(3.1))
                add_textbox(slide, title_data["title"], Inches(0.8), Inches(2.0),
                            Inches(11.5), Inches(1.2), font_size=44, bold=True,
                            color=WHITE, align=PP_ALIGN.LEFT)
                add_textbox(slide, "Generated by Nexa AI", Inches(0.8), Inches(3.3),
                            Inches(6), Inches(0.5), font_size=16, color=LIGHT_TXT)

                # â”€â”€ Content slides â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
                for sd in slides_data[1:]:
                    slide = prs.slides.add_slide(blank_layout)
                    add_bg(slide)

                    img_bytes = fetch_image(sd.get("image_query", topic))
                    if img_bytes:
                        from io import BytesIO as _BIO2
                        # Right-side image panel
                        pic = slide.shapes.add_picture(_BIO2(img_bytes),
                                                       Inches(7.5), Inches(0), Inches(5.83), H)
                        slide.shapes._spTree.remove(pic._element)
                        slide.shapes._spTree.insert(2, pic._element)
                        # Gradient overlay on image side
                        grad = slide.shapes.add_shape(1, Inches(6.5), Inches(0), Inches(6.83), H)
                        grad.fill.solid()
                        grad.fill.fore_color.rgb = DARK_BG
                        grad.line.fill.background()

                    # Accent bar
                    add_accent_bar(slide, top=Inches(1.1))

                    # Slide title
                    add_textbox(slide, sd.get("title", ""), Inches(0.5), Inches(0.2),
                                Inches(7.0), Inches(0.85), font_size=28, bold=True,
                                color=WHITE)

                    # Bullet points
                    bullets = sd.get("bullets", [])
                    y = Inches(1.35)
                    for bullet in bullets[:5]:
                        txb = slide.shapes.add_textbox(Inches(0.5), y, Inches(6.8), Inches(0.7))
                        tf  = txb.text_frame
                        tf.word_wrap = True
                        p   = tf.paragraphs[0]
                        run = p.add_run()
                        run.text = f"â€¢ {bullet}"
                        run.font.size = Pt(17)
                        run.font.color.rgb = LIGHT_TXT
                        y += Inches(0.75)

                    # Slide number
                    add_textbox(slide, str(prs.slides.index(slide) + 1),
                                Inches(12.5), Inches(7.1), Inches(0.5), Inches(0.3),
                                font_size=11, color=RGBColor(0x55, 0x55, 0x55))

                buffer = BytesIO()
                prs.save(buffer)
                buffer.seek(0)
                response = HttpResponse(
                    buffer.getvalue(),
                    content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
                )
                response['Content-Disposition'] = f'attachment; filename="{topic[:30]}.pptx"'
                return response
            
            else:
                response = HttpResponse(essay_text, content_type='text/plain')
                response['Content-Disposition'] = f'attachment; filename="{topic[:30]}.txt"'
                return response
                
        except json.JSONDecodeError:
            return JsonResponse({'error': 'Invalid JSON'}, status=400)
        except Exception as e:
            return JsonResponse({'error': str(e)}, status=500)
    
    return JsonResponse({'error': 'Method not allowed'}, status=405)


@login_required
def ai_material_action(request):
    """AJAX endpoint for AI actions on study materials (summarize, quiz, flashcards, explain)."""
    if request.method != 'POST':
        return JsonResponse({'error': 'Method not allowed'}, status=405)
    
    try:
        data = json.loads(request.body)
        action = data.get('action', '')
        material_title = data.get('material_title', 'Material')
        extracted_text = data.get('extracted_text', '')
        selected_model = resolve_model(data.get('model'))
        
        if not extracted_text:
            return JsonResponse({'success': False, 'error': 'No text content available to analyze.'})
        
        prompt = build_material_prompt(action, material_title, extracted_text, data)
        response = ask_ai(prompt, user=request.user, use_rag=False, model=selected_model)
        
        return JsonResponse({'success': True, 'response': response})
        
    except json.JSONDecodeError:
        return JsonResponse({'success': False, 'error': 'Invalid request data'}, status=400)
    except Exception as e:
        return JsonResponse({'success': False, 'error': str(e)}, status=500)


@login_required
def essay_editor(request, essay_id=None):
    """Render the manual essay editor page, optionally pre-loading an existing essay."""
    essay = None
    if essay_id:
        essay = get_object_or_404(EssayRequest, id=essay_id, user=request.user)
    return render(request, 'ai_tutor/essay_editor.html', {'essay': essay})


@login_required
def essay_autocomplete(request):
    """Return a short AI continuation suggestion for ghost-text autocomplete."""
    if request.method != 'POST':
        return JsonResponse({'error': 'POST required'}, status=405)
    try:
        data = json.loads(request.body)
        context = data.get('context', '').strip()
        if len(context) < 20:
            return JsonResponse({'ok': False, 'suggestion': ''})
        from .ai_utils import get_openai_client
        client = get_openai_client()
        resp = client.chat.completions.create(
            model=resolve_model(data.get('model')),
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You are an essay writing assistant. "
                        "Continue the text the user has written with 1-2 natural sentences. "
                        "Match their tone and style exactly. "
                        "Output ONLY the continuation â€” no quotes, no commentary."
                    )
                },
                {"role": "user", "content": context}
            ],
            max_tokens=80,
            temperature=0.7,
        )
        suggestion = resp.choices[0].message.content.strip()
        return JsonResponse({'ok': True, 'suggestion': suggestion})
    except Exception as e:
        return JsonResponse({'ok': False, 'error': str(e)})


@login_required
def essay_copilot(request):
    """Copilot actions: complete, expand, shorten, rephrase, fix, custom."""
    if request.method != 'POST':
        return JsonResponse({'error': 'POST required'}, status=405)
    try:
        data = json.loads(request.body)
        action = data.get('action', 'complete')
        selected_text = data.get('selected_text', '').strip()
        context = data.get('context', '').strip()
        custom_prompt = data.get('custom_prompt', '').strip()

        from .ai_utils import get_openai_client
        client = get_openai_client()

        prompts = {
            'complete': (
                "You are an expert essay writer. "
                "Continue writing from where the text ends. Write 2-4 paragraphs that flow naturally. "
                "Output ONLY the continuation.",
                context or selected_text
            ),
            'expand': (
                "You are an expert essay editor. "
                "Expand the following text with more detail, examples, and explanation. "
                "Keep the same tone. Output ONLY the expanded version.",
                selected_text or context
            ),
            'shorten': (
                "You are an expert essay editor. "
                "Shorten the following text while keeping all key ideas. "
                "Output ONLY the shortened version.",
                selected_text or context
            ),
            'rephrase': (
                "You are an expert essay editor. "
                "Rephrase the following text to sound more polished and varied. "
                "Output ONLY the rephrased version.",
                selected_text or context
            ),
            'fix': (
                "You are an expert proofreader. "
                "Fix all grammar, spelling, and punctuation errors in the following text. "
                "Output ONLY the corrected text.",
                selected_text or context
            ),
            'custom': (
                f"You are an expert essay writing assistant. {custom_prompt}. "
                "Output ONLY the result â€” no commentary.",
                selected_text or context
            ),
        }

        system_msg, user_content = prompts.get(action, prompts['complete'])
        if not user_content:
            return JsonResponse({'ok': False, 'error': 'No text to work with'})

        resp = client.chat.completions.create(
            model=resolve_model(data.get('model')),
            messages=[
                {"role": "system", "content": system_msg},
                {"role": "user", "content": user_content}
            ],
            max_tokens=800,
            temperature=0.7,
        )
        result = resp.choices[0].message.content.strip()
        return JsonResponse({'ok': True, 'result': result})
    except Exception as e:
        return JsonResponse({'ok': False, 'error': str(e)})


@login_required
def essay_editor_save(request, essay_id):
    """Save editor content (HTML) back to the essay record."""
    if request.method != 'POST':
        return JsonResponse({'error': 'POST required'}, status=405)
    essay = get_object_or_404(EssayRequest, id=essay_id, user=request.user)
    try:
        data = json.loads(request.body)
        content = data.get('content', '').strip()
        title = data.get('title', '').strip()
        if content:
            # Strip HTML tags for plain text storage
            import re
            plain = re.sub(r'<[^>]+>', ' ', content)
            plain = re.sub(r'\s+', ' ', plain).strip()
            essay.essay_text = plain
        if title:
            essay.topic = title
        essay.save()
        return JsonResponse({'ok': True})
    except Exception as e:
        return JsonResponse({'ok': False, 'error': str(e)})


def build_material_prompt(action, material_title, extracted_text, data=None):
    """Build the appropriate prompt based on the action."""
    
    text_preview = extracted_text[:3000]
    
    prompts = {
        'summarize': f"""You are an expert study assistant. Please provide a clear and comprehensive summary of the following content from "{material_title}". 

Focus on:
- Main ideas and key concepts
- Important details and supporting points
- Any definitions or formulas mentioned

Write in clear paragraphs, not bullet points. Make it suitable for study review.

Content:
{text_preview}

Summary:""",

        'quiz': f"""You are an expert study assistant. Create a quiz with 5 multiple choice questions based on the following content from "{material_title}".

For each question:
1. Make it test understanding, not just memorization
2. Provide 4 options (A, B, C, D)
3. Indicate the correct answer

Format your response like this:
Q1: [Your question here]
A. [Option A]
B. [Option B]
C. [Option C]
D. [Option D]
Answer: [Letter]

Content:
{text_preview}

Quiz:""",

        'flashcards': f"""You are an expert study assistant. Create 10 flashcards based on the following content from "{material_title}".

For each flashcard:
- Front: Key term or concept (brief, 1-3 words)
- Back: Definition or explanation (clear, concise)

Format your response as:
Term: Definition

Example:
Photosynthesis: The process by which plants convert light energy into chemical energy

Create at least 10 useful flashcards:

Content:
{text_preview}

Flashcards:""",

        'explain': f"""You are an expert teacher. Explain the key concepts from "{material_title}" in simple, easy-to-understand terms.

Your explanation should:
- Break down complex ideas into simpler parts
- Use analogies where helpful
- Be suitable for a student learning the topic for the first time
- Include examples where relevant

Write in clear paragraphs:

Content:
{text_preview}

Explanation:"""
    }
    
    if action == 'ask':
        user_message = data.get('message', '') if data else ''
        return f"""You are an expert study assistant helping with "{material_title}".

The student asks: {user_message}

Based on the following content, provide a helpful answer:

Content:
{text_preview}

Answer:"""
    
    return prompts.get(action, prompts['explain'])

