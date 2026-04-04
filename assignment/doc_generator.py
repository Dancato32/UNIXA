"""
Document generation utilities for Assignment AI.
Handles creation of Word, PowerPoint, and PDF files.
"""
import io
import re
import logging
import urllib.request
import urllib.parse

logger = logging.getLogger(__name__)


# ─── Helpers ────────────────────────────────────────────────────────────────

def _parse_slides(content):
    """
    Parse markdown-style content into a list of slide dicts:
    [{'title': str, 'bullets': [str], 'notes': str}]
    """
    slides = []
    current = None

    for raw_line in content.split('\n'):
        line = raw_line.strip()
        if not line:
            continue

        if line.startswith('## ') or line.startswith('# '):
            if current:
                slides.append(current)
            current = {'title': line.lstrip('#').strip(), 'bullets': [], 'notes': ''}
        elif line.startswith('- ') or line.startswith('* '):
            if current is None:
                current = {'title': 'Overview', 'bullets': [], 'notes': ''}
            current['bullets'].append(line.lstrip('-* ').strip())
        elif line.startswith('### '):
            if current is None:
                current = {'title': line.lstrip('#').strip(), 'bullets': [], 'notes': ''}
            else:
                current['bullets'].append(line.lstrip('#').strip())
        else:
            if current:
                if len(line) > 20:
                    current['notes'] += line + ' '
                else:
                    current['bullets'].append(line)

    if current:
        slides.append(current)

    return slides[:12]  # cap at 12 content slides


def _fetch_image_bytes(query, width=800, height=450):
    """
    Fetch a relevant AI generated image from Pollinations AI (no API key needed).
    Returns bytes or None on failure.
    """
    try:
        safe_query = urllib.parse.quote(query)
        url = f"https://image.pollinations.ai/prompt/{safe_query}?width={width}&height={height}&nologo=true"
        req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
        with urllib.request.urlopen(req, timeout=15) as resp:
            return resp.read()
    except Exception as e:
        logger.warning(f"Image fetch failed for '{query}': {e}")
        return None


# ─── Word ────────────────────────────────────────────────────────────────────

def generate_word_document(content, title):
    """Generate a styled Word document from content."""
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()

        # Page margins
        for section in doc.sections:
            section.top_margin = Inches(1)
            section.bottom_margin = Inches(1)
            section.left_margin = Inches(1.25)
            section.right_margin = Inches(1.25)

        # Title
        title_para = doc.add_heading(title, 0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in title_para.runs:
            run.font.size = Pt(24)

        doc.add_paragraph()

        for para in content.split('\n\n'):
            para = para.strip()
            if not para:
                continue
            if para.startswith('# '):
                doc.add_heading(para[2:], 1)
            elif para.startswith('## '):
                doc.add_heading(para[3:], 2)
            elif para.startswith('### '):
                doc.add_heading(para[4:], 3)
            elif re.match(r'^[-*] ', para):
                for item in para.split('\n'):
                    item = item.strip()
                    if item:
                        p = doc.add_paragraph(item.lstrip('-* '), style='List Bullet')
                        p.paragraph_format.space_after = Pt(4)
            else:
                p = doc.add_paragraph(para)
                p.paragraph_format.line_spacing = 1.5
                p.paragraph_format.space_after = Pt(8)

        return doc
    except Exception as e:
        raise Exception(f"Error generating Word document: {e}")


# ─── PowerPoint ──────────────────────────────────────────────────────────────

def generate_powerpoint_slides(content, title):
    """
    Generate a professional PowerPoint presentation with online images.
    Dark theme, large typography, one image per content slide.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches

        W = Inches(13.33)
        H = Inches(7.5)

        prs = Presentation()
        prs.slide_width = W
        prs.slide_height = H

        # ── Colour palette ──
        BG       = RGBColor(0x0a, 0x0a, 0x0a)
        SURFACE  = RGBColor(0x14, 0x14, 0x14)
        ACCENT   = RGBColor(0x63, 0x66, 0xf1)
        ACCENT_L = RGBColor(0xa5, 0xb4, 0xfc)
        WHITE    = RGBColor(0xff, 0xff, 0xff)
        GREY     = RGBColor(0x88, 0x88, 0x88)

        def fill_bg(slide, color=BG):
            from pptx.util import Pt
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = color

        def add_rect(slide, left, top, width, height, color, alpha=None):
            shape = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE
                left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.fill.background()
            return shape

        def add_textbox(slide, text, left, top, width, height,
                        size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
            from pptx.util import Pt
            txb = slide.shapes.add_textbox(left, top, width, height)
            tf = txb.text_frame
            tf.word_wrap = wrap
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            return txb

        # ── Title slide ──────────────────────────────────────────────────────
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        fill_bg(slide)

        # Accent bar left
        add_rect(slide, Inches(0), Inches(0), Inches(0.08), H, ACCENT)

        # Big title
        add_textbox(
            slide, title,
            Inches(0.5), Inches(2.2), Inches(8), Inches(2),
            size=44, bold=True, color=WHITE
        )
        # Subtitle
        add_textbox(
            slide, "Generated by NEXA AI",
            Inches(0.5), Inches(4.4), Inches(6), Inches(0.6),
            size=18, bold=False, color=GREY
        )
        # Accent dot
        dot = slide.shapes.add_shape(1, Inches(0.5), Inches(4.1), Inches(0.5), Inches(0.06))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()

        # ── Content slides ───────────────────────────────────────────────────
        if isinstance(content, list):
            parsed = content
        else:
            parsed = _parse_slides(content)

        for slide_data in parsed:
            slide_title = slide_data.get('title', 'Overview')
            bullets = slide_data.get('bullets', [])[:7]
            notes_text = slide_data.get('notes', '')
            image_prompt = slide_data.get('image_prompt', f"{title} {slide_title} realistic highly detailed")

            # Fetch image for this slide
            img_bytes = _fetch_image_bytes(image_prompt, 800, 450)

            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            fill_bg(slide)

            if img_bytes:
                # Right half: image with dark overlay
                img_stream = io.BytesIO(img_bytes)
                try:
                    pic = slide.shapes.add_picture(
                        img_stream,
                        Inches(7.5), Inches(0),
                        Inches(5.83), H
                    )
                    # Dark overlay on image
                    overlay = add_rect(
                        slide,
                        Inches(7.5), Inches(0),
                        Inches(5.83), H,
                        RGBColor(0x0a, 0x0a, 0x0a)
                    )
                    overlay.fill.solid()
                    overlay.fill.fore_color.rgb = RGBColor(0x0a, 0x0a, 0x0a)
                    # Set transparency via XML
                    from lxml import etree
                    solidFill = overlay.fill._xPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
                    if solidFill is not None:
                        srgb = solidFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                        if srgb is not None:
                            alpha_el = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                            alpha_el.set('val', '75000')  # 75% opacity overlay
                except Exception:
                    pass  # image failed, continue without it

            # Left accent bar
            add_rect(slide, Inches(0), Inches(0), Inches(0.06), H, ACCENT)

            # Slide number chip
            chip = add_rect(slide, Inches(0.35), Inches(0.3), Inches(0.5), Inches(0.28), SURFACE)

            # Title
            add_textbox(
                slide, slide_title,
                Inches(0.35), Inches(0.7), Inches(7), Inches(1.1),
                size=30, bold=True, color=WHITE
            )

            # Accent underline
            add_rect(slide, Inches(0.35), Inches(1.85), Inches(1.2), Inches(0.04), ACCENT)

            # Bullets
            if bullets:
                bullet_top = Inches(2.1)
                bullet_h = Inches(0.55)
                for i, bullet in enumerate(bullets):
                    # Bullet dot
                    dot = slide.shapes.add_shape(
                        1,
                        Inches(0.35), bullet_top + Inches(i * 0.62) + Inches(0.18),
                        Inches(0.1), Inches(0.1)
                    )
                    dot.fill.solid()
                    dot.fill.fore_color.rgb = ACCENT
                    dot.line.fill.background()

                    add_textbox(
                        slide, bullet,
                        Inches(0.6), bullet_top + Inches(i * 0.62),
                        Inches(6.7), bullet_h,
                        size=16, bold=False, color=RGBColor(0xdd, 0xdd, 0xdd)
                    )

            # Speaker notes
            if notes_text:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes_text.strip()

        # ── Thank you slide ──────────────────────────────────────────────────
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        fill_bg(slide)
        add_rect(slide, Inches(0), Inches(0), Inches(0.08), H, ACCENT)
        add_textbox(
            slide, "Thank You",
            Inches(0.5), Inches(2.8), Inches(8), Inches(1.5),
            size=48, bold=True, color=WHITE, align=PP_ALIGN.LEFT
        )
        add_textbox(
            slide, "Generated by NEXA AI Assignment Assistant",
            Inches(0.5), Inches(4.5), Inches(8), Inches(0.6),
            size=16, bold=False, color=GREY
        )

        return prs

    except Exception as e:
        raise Exception(f"Error generating PowerPoint: {e}")


# ─── PDF ─────────────────────────────────────────────────────────────────────

def generate_pdf_document(content, title):
    """Generate a styled PDF document from content."""
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.enums import TA_LEFT, TA_CENTER
        from reportlab.lib.colors import HexColor

        buffer = io.BytesIO()
        doc = SimpleDocTemplate(
            buffer, pagesize=letter,
            topMargin=0.75 * inch, bottomMargin=0.75 * inch,
            leftMargin=1.25 * inch, rightMargin=1.25 * inch
        )

        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'NexaTitle', parent=styles['Heading1'],
            fontSize=26, spaceAfter=24, alignment=TA_CENTER,
            textColor=HexColor('#1a1a2e')
        )
        h1_style = ParagraphStyle(
            'NexaH1', parent=styles['Heading2'],
            fontSize=18, spaceAfter=10, spaceBefore=18,
            textColor=HexColor('#6366f1')
        )
        h2_style = ParagraphStyle(
            'NexaH2', parent=styles['Heading3'],
            fontSize=14, spaceAfter=8, spaceBefore=14,
            textColor=HexColor('#4f52d9')
        )
        body_style = ParagraphStyle(
            'NexaBody', parent=styles['Normal'],
            fontSize=11, spaceAfter=10, leading=18, alignment=TA_LEFT
        )
        bullet_style = ParagraphStyle(
            'NexaBullet', parent=styles['Normal'],
            fontSize=11, spaceAfter=5, leading=16,
            leftIndent=20, bulletIndent=8
        )

        story = [Paragraph(title, title_style), Spacer(1, 0.2 * inch)]

        for para in content.split('\n\n'):
            para = para.strip()
            if not para:
                continue
            if para.startswith('# '):
                story.append(Paragraph(para[2:], h1_style))
            elif para.startswith('## '):
                story.append(Paragraph(para[3:], h1_style))
            elif para.startswith('### '):
                story.append(Paragraph(para[4:], h2_style))
            elif re.match(r'^[-*] ', para):
                for item in para.split('\n'):
                    item = item.strip().lstrip('-* ')
                    if item:
                        story.append(Paragraph(f"• {item}", bullet_style))
            else:
                clean = para.replace('\n', ' ').replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                story.append(Paragraph(clean, body_style))
                story.append(Spacer(1, 4))

        doc.build(story)
        return buffer

    except Exception as e:
        raise Exception(f"Error generating PDF: {e}")


# ─── Text extractor ──────────────────────────────────────────────────────────

def extract_text_from_file(file):
    """Extract text from uploaded file based on file type."""
    text = ""
    file_name = file.name.lower()

    try:
        if file_name.endswith('.pdf'):
            import PyPDF2
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() or ""

        elif file_name.endswith('.docx'):
            from docx import Document
            doc = Document(file)
            for para in doc.paragraphs:
                text += para.text + "\n"

        elif file_name.endswith('.pptx') or file_name.endswith('.ppt'):
            from pptx import Presentation
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

        elif file_name.endswith('.txt'):
            text = file.read().decode('utf-8', errors='ignore')

        elif file_name.lower().endswith(('.jpg', '.jpeg', '.png', '.webp', '.gif')):
            text = "[Image file — text will be extracted via vision AI]"

        return text
    except Exception as e:
        return f"[Could not extract text from file: {e}]"
