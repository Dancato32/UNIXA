from nexa.rate_limit import ai_rate_limit, upload_rate_limit
from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth.decorators import login_required
from django.contrib import messages
from django.core.paginator import Paginator
from django.http import JsonResponse, FileResponse, StreamingHttpResponse, HttpResponse
from django.conf import settings
from .models import StudyMaterial, SavedFlashcardDeck, SavedPodcast, SavedStudySong, ConceptNode, StudentConceptState
from .forms import StudyMaterialForm
from ai_tutor.ai_utils import ask_ai, resolve_model
from ai_tutor.models import ChatThread, Conversation
import os
import json
import uuid
import re


def extract_text_from_file(file_path, file_extension):
    """
    Extract text from uploaded files.
    Supports: .docx (Word), .pptx (PowerPoint), .pdf
    """
    text = ""
    
    if file_extension == '.docx':
        # Extract text from Word documents
        try:
            import docx
            doc = docx.Document(file_path)
            for paragraph in doc.paragraphs:
                text += paragraph.text + '\n'
        except ImportError:
            text = "[Word text extraction requires 'python-docx'. Install it with: pip install python-docx]"
        except Exception as e:
            text = f"[Error extracting Word content: {str(e)}]"
            
    elif file_extension == '.pptx':
        # Extract text from PowerPoint presentations
        try:
            import pptx
            prs = pptx.Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + '\n'
        except ImportError:
            text = "[PowerPoint text extraction requires 'python-pptx'. Install it with: pip install python-pptx]"
        except Exception as e:
            text = f"[Error extracting PowerPoint content: {str(e)}]"
            
    elif file_extension == '.pdf':
        # Simple PDF text extraction using PyPDF2
        try:
            import PyPDF2
            with open(file_path, 'rb') as pdf_file:
                pdf_reader = PyPDF2.PdfReader(pdf_file)
                for page in pdf_reader.pages:
                    text += page.extract_text() or ''
        except ImportError:
            text = "[PDF text extraction requires 'PyPDF2'. Install it with: pip install PyPDF2]"
        except Exception as e:
            text = f"[Error extracting PDF content: {str(e)}]"
    
    return text


def extract_text_from_memory(uploaded_file, file_extension):
    """Extract text from an in-memory uploaded file (used with cloud storage)."""
    import tempfile, shutil
    text = ""
    try:
        uploaded_file.seek(0)
        suffix = file_extension or '.tmp'
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            shutil.copyfileobj(uploaded_file, tmp)
            tmp_path = tmp.name
        text = extract_text_from_file(tmp_path, file_extension)
        os.remove(tmp_path)
    except Exception as e:
        text = f"[Error extracting text: {str(e)}]"
    return text


@login_required
def upload_material(request):
    """Allow logged-in students to upload study materials."""
    if request.method == 'POST':
        form = StudyMaterialForm(request.POST, request.FILES)
        if form.is_valid():
            uploaded_file = request.FILES.get('file')
            MAX_SIZE = 200 * 1024 * 1024  # 200 MB â€” local filesystem, no Cloudinary limit
            if uploaded_file and uploaded_file.size > MAX_SIZE:
                form.add_error('file', f'File too large ({uploaded_file.size // (1024*1024)} MB). Maximum allowed size is 200 MB.')
                return render(request, 'materials/upload.html', {'form': form, 'title': 'Upload Study Material'})
            material = form.save(commit=False)
            material.owner = request.user
            try:
                material.save()
            except Exception as e:
                err = str(e)
                if 'size' in err.lower() or 'large' in err.lower() or '10485760' in err:
                    form.add_error('file', 'File too large. Maximum allowed size is 200 MB.')
                else:
                    form.add_error(None, f'Upload failed: {err}')
                return render(request, 'materials/upload.html', {'form': form, 'title': 'Upload Study Material'})
            file_extension = os.path.splitext(uploaded_file.name)[1].lower() if uploaded_file else ''
            try:
                file_path = material.file.path
                extracted_text = extract_text_from_file(file_path, file_extension)
            except (NotImplementedError, ValueError, AttributeError):
                extracted_text = extract_text_from_memory(uploaded_file, file_extension)
            material.extracted_text = extracted_text
            material.save(update_fields=['extracted_text'])
            if extracted_text.startswith('[') and 'requires' in extracted_text:
                messages.warning(request, f'Material uploaded, but {extracted_text}')
            else:
                messages.success(request, 'Material uploaded successfully!')
            return redirect('list_materials')
    else:
        form = StudyMaterialForm()
    return render(request, 'materials/upload.html', {'form': form, 'title': 'Upload Study Material'})


@login_required
def upload_material_ajax(request):
    """AJAX upload endpoint — auto-detects material_type from file extension."""
    if request.method != 'POST':
        return JsonResponse({'error': 'POST required'}, status=405)

    uploaded_file = request.FILES.get('file')
    if not uploaded_file:
        return JsonResponse({'success': False, 'errors': {'file': ['No file was submitted.']}}, status=400)

    # Auto-detect material_type from extension
    ext = os.path.splitext(uploaded_file.name)[1].lower()
    ext_to_type = {
        '.pdf':  'PDF',
        '.docx': 'Word',
        '.doc':  'Word',
        '.pptx': 'PowerPoint',
        '.ppt':  'PowerPoint',
    }
    material_type = ext_to_type.get(ext)
    if not material_type:
        return JsonResponse({
            'success': False,
            'errors': {'file': [f'Unsupported file type "{ext}". Please upload PDF, Word (.docx), or PowerPoint (.pptx).']}
        }, status=400)

    # Build POST data with auto-detected type and title fallback
    post_data = request.POST.copy()
    post_data['material_type'] = material_type
    if not post_data.get('title', '').strip():
        post_data['title'] = os.path.splitext(uploaded_file.name)[0][:255]

    form = StudyMaterialForm(post_data, request.FILES)
    if not form.is_valid():
        return JsonResponse({'success': False, 'errors': form.errors}, status=400)

    MAX_SIZE = 200 * 1024 * 1024  # 200 MB
    if uploaded_file.size > MAX_SIZE:
        return JsonResponse({
            'success': False,
            'errors': {'file': [f'File too large ({uploaded_file.size // (1024*1024)} MB). Maximum is 200 MB.']}
        }, status=400)

    material = form.save(commit=False)
    material.owner = request.user
    material.material_type = material_type
    try:
        material.save()
    except Exception as e:
        err = str(e)
        if 'size' in err.lower() or 'large' in err.lower() or '10485760' in err:
            return JsonResponse({'success': False, 'errors': {'file': ['File too large. Maximum is 200 MB.']}}, status=400)
        return JsonResponse({'success': False, 'errors': {'__all__': [f'Upload failed: {err}']}}, status=500)

    # Extract text for AI features
    try:
        file_path = material.file.path
        extracted_text = extract_text_from_file(file_path, ext)
    except (NotImplementedError, ValueError, AttributeError):
        extracted_text = extract_text_from_memory(uploaded_file, ext)

    material.extracted_text = extracted_text
    material.save(update_fields=['extracted_text'])
    return JsonResponse({'success': True, 'redirect': '/materials/list/'})



@login_required
def list_materials(request):
    """List all uploaded materials by the logged-in student."""
    
    materials_list = StudyMaterial.objects.filter(owner=request.user)
    
    # Add pagination
    paginator = Paginator(materials_list, 10)  # Show 10 materials per page
    page_number = request.GET.get('page')
    materials = paginator.get_page(page_number)
    
    return render(request, 'materials/list.html', {
        'materials': materials,
        'title': 'My Study Materials'
    })


@login_required
def material_detail(request, pk):
    """View details of a specific study material."""
    
    material = get_object_or_404(StudyMaterial, pk=pk, owner=request.user)
    return render(request, 'materials/detail.html', {
        'material': material,
        'title': material.title
    })


@login_required
def delete_material(request, pk):
    """Delete a study material."""
    
    material = get_object_or_404(StudyMaterial, pk=pk, owner=request.user)
    
    if request.method == 'POST':
        material.delete()
        messages.success(request, 'Material deleted successfully.')
        return redirect('list_materials')
    
    return render(request, 'materials/delete.html', {
        'material': material,
        'title': 'Delete Material'
    })


@login_required
def podcast_view(request, pk):
    """View for the Podcast Learning Mode."""
    material = get_object_or_404(StudyMaterial, pk=pk, owner=request.user)
    
    return render(request, 'materials/podcast.html', {
        'material': material,
        'title': f'Podcast: {material.title}'
    })


@login_required
def generate_podcast_ajax(request):
    """AJAX endpoint to generate podcast script and ElevenLabs audio."""
    if request.method != 'POST':
        return JsonResponse({'error': 'POST request required'}, status=405)

    try:
        data = json.loads(request.body)
        material_id = data.get('material_id')

        material = get_object_or_404(StudyMaterial, pk=material_id, owner=request.user)

        if not material.extracted_text:
            return JsonResponse({'error': 'No text content available for this material'}, status=400)

        text_content = material.extracted_text[:4000]
        selected_model = resolve_model(data.get('model'))
        prompt = build_podcast_prompt(text_content, material.title)
        podcast_script = ask_ai(prompt, user=request.user, use_rag=False, model=selected_model)

        word_count = len(podcast_script.split())
        duration_mins = round(word_count / 130)

        # Generate Resemble.ai audio â€” returns base64 data URL directly
        audio_url = None
        print(f"[PODCAST] Generating Resemble audio for material {material_id}, script length: {len(podcast_script)}")
        audio_data_url = generate_podcast_audio_elevenlabs(podcast_script, material_id)
        print(f"[PODCAST] audio_data_url result: {bool(audio_data_url)}")
        if audio_data_url:
            audio_url = audio_data_url
        print(f"[PODCAST] audio_url set: {bool(audio_url)}")

        return JsonResponse({
            'success': True,
            'raw_content': podcast_script,
            'material_title': material.title,
            'word_count': word_count,
            'duration_estimate': f'~{duration_mins} minutes',
            'audio_url': audio_url,
        })

    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


def generate_podcast_audio(script_text, material_id):
    """Generate podcast audio - now directly uses OpenAI TTS."""
    # OpenRouter doesn't support TTS, so we go straight to OpenAI
    return generate_podcast_audio_direct(script_text, material_id)


def generate_podcast_audio_direct(script_text, material_id):
    """Generate podcast audio using direct OpenAI TTS API."""
    try:
        from openai import OpenAI
        import os
        
        # Get direct OpenAI API key only
        api_key = os.getenv("OPENAI_API_KEY")
        
        if not api_key:
            print("No direct OpenAI API key for TTS")
            return None
        
        # Create OpenAI client
        client = OpenAI(api_key=api_key)
        
        # Prepare full text for audio (combine all segments)
        full_text = script_text
        if isinstance(script_text, dict):
            full_text = " ".join([seg['content'] for seg in script_text.get('segments', [])])
        
        # Limit text length for TTS
        full_text = full_text[:2500]
        
        if not full_text:
            return None
        
        # Create audio directory if not exists
        audio_dir = os.path.join(settings.MEDIA_ROOT, 'podcasts')
        os.makedirs(audio_dir, exist_ok=True)
        
        # Generate unique filename
        filename = f"podcast_{material_id}_{uuid.uuid4().hex[:8]}.mp3"
        filepath = os.path.join(audio_dir, filename)
        
        # Generate speech using OpenAI's TTS
        response = client.audio.speech.create(
            model="tts-1",
            voice="nova",
            input=full_text,
            response_format="mp3"
        )
        
        # Save audio file
        response.stream_to_file(filepath)
        
        return filename
        
    except Exception as e:
        print(f"Direct TTS Error: {e}")
        return None


def _get_resemble_voice_uuid(api_key):
    """Returns the NEXA custom voice UUID."""
    return "f644f59c"  # My Custom Voice NEXA


def generate_podcast_audio_elevenlabs(script_text, material_id):
    """Generate podcast audio using Resemble.ai TTS API. Returns base64 data URL (no file system needed)."""
    try:
        import os
        import requests

        api_key = os.environ.get("RESEMBLE_API_KEY") or os.getenv("RESEMBLE_API_KEY") or getattr(settings, 'RESEMBLE_API_KEY', None)
        print(f"[RESEMBLE] api_key present: {bool(api_key)}, value starts: {api_key[:8] if api_key else 'NONE'}")
        if not api_key:
            print("[RESEMBLE] No Resemble API key found")
            return None

        full_text = script_text
        if isinstance(script_text, dict):
            full_text = " ".join([seg['content'] for seg in script_text.get('segments', [])])
        full_text = full_text[:2000]
        print(f"[RESEMBLE] text length after truncation: {len(full_text)}")
        if not full_text:
            return None

        voice_uuid = _get_resemble_voice_uuid(api_key)
        print(f"[RESEMBLE] using voice_uuid: {voice_uuid}")

        print(f"[RESEMBLE] calling synthesize API...")
        response = requests.post(
            "https://f.cluster.resemble.ai/synthesize",
            headers={"Authorization": "Bearer " + api_key, "Content-Type": "application/json"},
            json={"voice_uuid": voice_uuid, "data": full_text, "output_format": "mp3"},
            timeout=120
        )

        print(f"[RESEMBLE] response status: {response.status_code}")
        if response.status_code == 200:
            audio_b64 = response.json().get("audio_content")
            print(f"[RESEMBLE] audio_content present: {bool(audio_b64)}")
            if audio_b64:
                # Return as base64 data URL â€” no file system, works on iOS/Android
                data_url = f"data:audio/mpeg;base64,{audio_b64}"
                print(f"[RESEMBLE] returning data URL (length: {len(data_url)})")
                return data_url
        print(f"[RESEMBLE] API Error {response.status_code}: {response.text[:300]}")
        return None

    except Exception as e:
        print(f"[RESEMBLE] Exception: {e}")
        import traceback
        traceback.print_exc()
        return None


def generate_podcast_audio_streaming(request, script_text, material_id):
    """Generate podcast audio using streaming TTS for immediate playback."""
    try:
        from openai import OpenAI
        import os

        # TTS requires a direct OpenAI key â€” OpenRouter does NOT support TTS
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            print("No OPENAI_API_KEY for TTS â€” skipping streaming audio")
            return None, None

        client = OpenAI(api_key=api_key)
        
        # Prepare full text
        full_text = script_text
        if isinstance(script_text, dict):
            full_text = " ".join([seg['content'] for seg in script_text.get('segments', [])])
        
        full_text = full_text[:2500]
        
        if not full_text:
            return None, None
        
        # Create audio directory
        audio_dir = os.path.join(settings.MEDIA_ROOT, 'podcasts')
        os.makedirs(audio_dir, exist_ok=True)
        
        # Generate unique filename
        filename = f"podcast_{material_id}_{uuid.uuid4().hex[:8]}.mp3"
        filepath = os.path.join(audio_dir, filename)
        
        # Generate with streaming for immediate playback
        response = client.audio.speech.create(
            model="tts-1-hd",  # Higher quality TTS
            voice="nova",
            input=full_text,
            response_format="mp3"
        )
        
        # Save while generating (streams to file)
        response.stream_to_file(filepath)
        
        return filename, filepath
        
    except Exception as e:
        print(f"Streaming TTS Error: {e}")
        return None, None


@login_required
def serve_podcast_audio(request, material_id, filename=None):
    """Serve the generated podcast audio file with HTTP Range support for browser <audio>."""
    try:
        audio_dir = os.path.join(settings.MEDIA_ROOT, 'podcasts')

        if not os.path.exists(audio_dir):
            return HttpResponse('Audio not found', status=404)

        # Resolve filepath
        filepath = None
        if filename:
            candidate = os.path.join(audio_dir, filename)
            if os.path.exists(candidate):
                filepath = candidate
            else:
                matches = [f for f in os.listdir(audio_dir)
                           if f.startswith(f'podcast_{material_id}_') and filename in f]
                if matches:
                    filepath = os.path.join(audio_dir, matches[0])

        if not filepath:
            files = [f for f in os.listdir(audio_dir) if f.startswith(f'podcast_{material_id}_')]
            if not files:
                return HttpResponse('Audio not found', status=404)
            latest = max(files, key=lambda f: os.path.getctime(os.path.join(audio_dir, f)))
            filepath = os.path.join(audio_dir, latest)

        file_size = os.path.getsize(filepath)
        content_type = 'audio/mpeg'

        # Handle HTTP Range requests (required for browser audio seeking/streaming)
        range_header = request.META.get('HTTP_RANGE', '').strip()
        if range_header:
            range_match = range_header.replace('bytes=', '').split('-')
            first_byte = int(range_match[0]) if range_match[0] else 0
            last_byte = int(range_match[1]) if range_match[1] else file_size - 1
            last_byte = min(last_byte, file_size - 1)
            length = last_byte - first_byte + 1

            with open(filepath, 'rb') as f:
                f.seek(first_byte)
                data = f.read(length)

            response = HttpResponse(data, status=206, content_type=content_type)
            response['Content-Range'] = f'bytes {first_byte}-{last_byte}/{file_size}'
            response['Accept-Ranges'] = 'bytes'
            response['Content-Length'] = str(length)
            return response

        # Full file response
        response = FileResponse(open(filepath, 'rb'), content_type=content_type)
        response['Content-Length'] = str(file_size)
        response['Accept-Ranges'] = 'bytes'
        return response

    except Exception as e:
        return HttpResponse(str(e), status=500)


def build_podcast_prompt(text_content, title):
    """Build a high-fidelity script for 'The Alex & Sam Masterclass Podcast'."""
    return f"""You are a world-class Podcast Producer directing the "Alex & Sam Masterclass". Your goal is to write a highly natural, engaging, and premium podcast script.
    
TOPIC: {title}

THE HOSTS:
- ALEX (Male): The expert. Professional, friendly, deep-thinking, and clear.
- SAM (Female): The color-commentator. Insightful, relatable, asking the right questions, and connecting ideas.

STRICT RULES FOR AUDIO REALISM:
1. START EXACTLY LIKE THIS:
   Alex: Hi, my name is Alex.
   Sam: And my name is Sam. Today we're going to be talking about {title}.
2. ZERO ROBOTICS: After the intro, transition into a highly natural conversation. Treat it like a natural discussion, not constantly saying each other's names. Start with brief casual banter (like grabbing coffee or their weekend) before transitioning to the subject.
3. PURE DIALOGUE ONLY: Write only what they SAY. No stage directions [Laughs].
4. NO BRACKETS: No bracketed text at all.
5. NATURAL OVERLAPS: Use natural conversation fillers ("Um", "Hmm", "Right", "Exactly", "Wait") fluidly to make them sound like they are sitting in the same room.

STRUCTURE:
- INTRO: [Exact intro -> Brief natural banter -> Smooth transition]
- DECONSTRUCTION: [Breaking down the first major concept naturally]
- THE "AHA" MOMENT: [A back-and-forth epiphany about the hardest part of the topic]
- OUTRO: [A quick, confident sign-off]

STUDY MATERIAL:
{text_content}

Now, start the session:"""

def parse_podcast_response(ai_response, title):
    """Parse the AI response into structured podcast segments."""
    import re
    
    segments = []
    
    # Default structure
    default_segments = [
        {'type': 'intro', 'title': 'Welcome to Your Study Session', 'content': f"Naruto: Believe it! We're diving into {title} today! Sasuke: Quiet, Naruto. Focus on the lesson.", 'visual': 'intro'},
        {'type': 'concept', 'title': 'Getting Started', 'content': 'Loading your personalized training content...', 'visual': 'loading'},
    ]
    
    # Try to parse the AI response
    try:
        # Split by segment markers
        segment_pattern = r'-(INTRO|CONCEPT_\d+|EXAMPLE|SUMMARY|OUTRO):\s*'
        parts = re.split(segment_pattern, ai_response)
        
        if len(parts) > 1:
            # We have segmented content
            i = 1
            while i < len(parts):
                marker = parts[i].strip() if i < len(parts) else ''
                content = parts[i+1].strip() if i+1 < len(parts) else ''
                
                if content:
                    segment_type = 'intro' if 'INTRO' in marker.upper() else 'concept'
                    if 'EXAMPLE' in marker.upper():
                        segment_type = 'example'
                    elif 'SUMMARY' in marker.upper():
                        segment_type = 'summary'
                    elif 'OUTRO' in marker.upper():
                        segment_type = 'outro'
                    
                    # Extract title from first line if possible
                    lines = content.split('\n')
                    segment_title = lines[0].strip() if lines else marker.replace('_', ' ').title()
                    segment_content = ' '.join(lines[1:]) if len(lines) > 1 else content
                    
                    segments.append({
                        'type': segment_type,
                        'title': segment_title[:50],
                        'content': segment_content[:500],
                        'visual': segment_type
                    })
                
                i += 2
        
        # If no segments parsed, create from raw content
        if not segments:
            # Split by paragraphs
            paragraphs = ai_response.split('\n\n')
            for i, para in enumerate(paragraphs[:8]):
                if para.strip() and len(para.strip()) > 20:
                    seg_type = 'intro' if i == 0 else 'summary' if i >= len(paragraphs) - 2 else 'concept'
                    segments.append({
                        'type': seg_type,
                        'title': f"Part {i+1}" if seg_type == 'concept' else ('Introduction' if seg_type == 'intro' else 'Summary'),
                        'content': para.strip()[:500],
                        'visual': seg_type
                    })
    except Exception as e:
        # Fallback to default
        pass
    
    # Ensure we have at least some content
    if not segments:
        # Create basic segments from raw response
        content_lines = ai_response.split('\n')
        meaningful_lines = [l.strip() for l in content_lines if len(l.strip()) > 30][:10]
        
        for i, line in enumerate(meaningful_lines):
            seg_type = 'intro' if i == 0 else 'summary' if i >= len(meaningful_lines) - 1 else 'concept'
            segments.append({
                'type': seg_type,
                'title': f'Section {i+1}',
                'content': line,
                'visual': seg_type
            })
    
    # Ensure minimum segments
    if len(segments) < 3:
        segments = [
            {'type': 'intro', 'title': 'Welcome', 'content': f"Alex: Hi, my name is Alex. Sam: And my name is Sam. Today we're going to be talking about {title}.", 'visual': 'intro'}
        ] + segments
    
    return {
        'title': f"Alex & Sam: {title}",
        'segments': segments[:10]
    }


@login_required
def materials_count_api(request):
    """API endpoint to get count of user's study materials."""
    count = StudyMaterial.objects.filter(owner=request.user).count()
    return JsonResponse({'count': count})


@login_required
def podcast_question_ajax(request):
    """AJAX endpoint to answer questions about podcast content."""
    if request.method != 'POST':
        return JsonResponse({'error': 'POST request required'}, status=405)
    
    try:
        data = json.loads(request.body)
        question = data.get('question', '').strip()
        podcast_context = data.get('podcast_context', '')
        material_id = data.get('material_id')
        
        if not question:
            return JsonResponse({'error': 'Question is required'}, status=400)
        
        # Build AI prompt with podcast context
        prompt = f"""You are Nexa, an AI tutor helping a student understand their study material. 

The student is listening to a podcast about their study material and has paused to ask a question.

Podcast Content:
{podcast_context[:2000]}

Student's Question: {question}

Provide a clear, concise answer (2-3 sentences) that directly addresses their question. Be friendly and encouraging."""
        
        # Get AI answer
        selected_model = resolve_model(data.get('model'))
        answer = ask_ai(prompt, user=request.user, use_rag=False, model=selected_model)
        
        # Generate audio for the answer natively using Edge TTS (matching the podcast hosts)
        audio_url = None
        try:
            audio_filename = generate_answer_audio_host(answer, material_id)
            if audio_filename:
                audio_url = f'/materials/podcast/answer-audio/{material_id}/{audio_filename}'
        except Exception as e:
            print(f"Host TTS failed: {e}")
        
        return JsonResponse({
            'success': True,
            'answer': answer,
            'audio_url': audio_url
        })
        
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


def generate_answer_audio_host(answer_text, material_id):
    """Generate answer audio natively using Edge TTS (randomly alternating between podcast hosts)"""
    try:
        import os
        import uuid
        import random
        import asyncio
        import concurrent.futures
        import edge_tts
        from django.conf import settings

        text = answer_text[:1000]
        if not text:
            return None

        audio_dir = os.path.join(settings.MEDIA_ROOT, 'podcast_answers')
        os.makedirs(audio_dir, exist_ok=True)

        filename = f"answer_{material_id}_{uuid.uuid4().hex[:8]}.mp3"
        filepath = os.path.join(audio_dir, filename)

        # Randomly choose between Alex's and Sam's actual Edge TTS voice
        voice = random.choice(['en-US-AndrewNeural', 'en-US-AvaNeural'])

        async def _bake_answer():
            await edge_tts.Communicate(text, voice).save(filepath)

        # Use a dedicated thread with a fresh event loop — safe on Gunicorn/production
        def _runner():
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
            try:
                loop.run_until_complete(_bake_answer())
            finally:
                loop.close()

        with concurrent.futures.ThreadPoolExecutor(max_workers=1) as pool:
            pool.submit(_runner).result()

        return filename
    except Exception as e:
        print(f"Host answer TTS Exception: {e}")
        return None


@login_required
def serve_answer_audio(request, material_id, filename):
    """Serve the generated answer audio file."""
    try:
        audio_dir = os.path.join(settings.MEDIA_ROOT, 'podcast_answers')
        filepath = os.path.join(audio_dir, filename)
        
        if os.path.exists(filepath):
            return FileResponse(open(filepath, 'rb'), content_type='audio/mpeg')
        
        return JsonResponse({'error': 'Audio file not found'}, status=404)
        
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


@login_required
def summarize_view(request, pk):
    """Dedicated page: AI summary of a study material."""
    material = get_object_or_404(StudyMaterial, pk=pk, owner=request.user)
    return render(request, 'materials/summarize.html', {
        'material': material, 'title': f'Summary: {material.title}'
    })


@login_required
def summarize_ajax(request, pk):
    """AJAX: generate summary."""
    if request.method != 'POST':
        return JsonResponse({'error': 'POST required'}, status=405)
    material = get_object_or_404(StudyMaterial, pk=pk, owner=request.user)
    if not material.extracted_text:
        return JsonResponse({'error': 'No text could be extracted from this material. Try re-uploading it.'}, status=400)
    prompt = f"""Summarize the following study material titled "{material.title}" in clear, structured paragraphs.
Cover all main ideas, key concepts, and important details. Write for a student reviewing for an exam.
Do not use markdown bold (**). Use plain text only.

Content:
{material.extracted_text[:4000]}

Summary:"""
    try:
        data = {}
        if request.body:
            try: data = json.loads(request.body)
            except: pass
        selected_model = resolve_model(data.get('model'))
        result = ask_ai(prompt, user=request.user, use_rag=False, model=selected_model)
        return JsonResponse({'success': True, 'result': result})
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


@login_required
def quiz_view(request, pk):
    """Dedicated page: AI-generated interactive quiz with meta-cognitive awareness."""
    material = get_object_or_404(StudyMaterial, pk=pk, owner=request.user)
    concepts = material.concepts.all()
    from .models import StudentConceptState
    
    concept_data = []
    for c in concepts:
        has_risk = False
        try:
            state = StudentConceptState.objects.get(user=request.user, concept=c)
            if (state.error_profile or {}).get("misconception", 0) >= 2:
                has_risk = True
        except StudentConceptState.DoesNotExist:
            pass
        
        concept_data.append({
            "id": c.id,
            "name": c.name,
            "definition": c.definition,
            "recurrence_risk": has_risk
        })
    
    return render(request, "materials/quiz.html", {
        "material": material,
        "title": f"Quiz: {material.title}",
        "concepts": concept_data
    })

@login_required
def quiz_ajax(request, pk):
    """AJAX: generate quiz with pattern-based rewiring and persistence escalation."""
    if request.method != "POST":
        return JsonResponse({"error": "POST required"}, status=405)
    material = get_object_or_404(StudyMaterial, pk=pk, owner=request.user)
    if not material.extracted_text:
        return JsonResponse({"error": "No text could be extracted from this material."}, status=400)
    
    data = {}
    if request.body:
        try: data = json.loads(request.body)
        except: pass
    
    selected_topics = data.get("topics", [])
    custom_topic = data.get("custom_topic", "").strip()
    difficulty = data.get("difficulty", "academic")
    selected_model = resolve_model(data.get("model"))
    
    from .models import ConceptNode, StudentConceptState
    from django.utils import timezone
    import datetime
    week_ago = timezone.now() - datetime.timedelta(days=7)
    
    # Context building & Persistence Detection
    context_info = ""
    dep_info = ""
    escalation_instructions = ""
    
    if selected_topics:
        topic_objs = material.concepts.filter(name__in=selected_topics)
        for t in topic_objs:
            context_info += f"- {t.name}: {t.definition}\n"
            # Detect recurrence
            try:
                state = StudentConceptState.objects.get(user=request.user, concept=t)
                errors = state.error_profile or {}
                misconception_count = errors.get("misconception", 0)
                if misconception_count >= 2:
                    escalation_instructions += f"- ESCALATE {t.name}: The student has a PERSISTANT MISCONCEPTION on this topic. Instead of a standard fix, use an ANALOGY or FIRST PRINCIPLES breakdown in the remediation bridge.\n"
            except StudentConceptState.DoesNotExist:
                pass

            for prereq in t.prerequisites.all():
                dep_info += f"  - {t.name} depends on: {prereq.name}\n"
    
    if custom_topic:
        context_info += f"- FOCUS AREA: {custom_topic}\n"
    
    target_desc = "the entire material" if not context_info else "the specific topics listed below"
    
    diff_instructions = {
        "casual": "Focus on RECALL and basic DEFINITIONS. Questions should test whether the student remembers key terms and facts. Keep language simple.",
        "academic": "Focus on APPLICATION and WHY questions. Test understanding of how concepts work and connect to each other.",
        "mastery": "Focus on SYNTHESIS and REASONING. Ask multi-concept questions that require deep analysis, comparison, or evaluation."
    }
    diff_prompt = diff_instructions.get(difficulty, diff_instructions["academic"])
    
    prompt = f"""You are a cognitive science assessment expert. Create a 5-question multiple choice quiz.

DIFFICULTY LEVEL: {difficulty.upper()}
{diff_prompt}

{"TARGET TOPICS:\\n" + context_info if context_info else "Cover the entire material broadly."}
{"PREREQUISITE MAP:\\n" + dep_info if dep_info else ""}
{"URGENT ESCALATIONS:\\n" + escalation_instructions if escalation_instructions else ""}

MATERIAL TEXT (Excerpt):
{material.extracted_text[:4000]}

Return ONLY a valid JSON array. NO markdown, NO code blocks, NO extra text.
Each object MUST have: "q", "concept", "opts", "ans", "explanation", "traps", "remediation", "dependencies".
For "traps": map letter to {{"error_type", "trap_explanation"}}. Error Types: misconception | partial_confusion | misapplied_rule | calculation | recall | guessing.
For "remediation": map to {{"bridge_question", "bridge_options", "bridge_answer", "bridge_explanation"}}. 
IMPORTANT: If an ESCALATION is requested for a topic, the "bridge_explanation" MUST be an analogy or first-principles breakdown.

JSON:"""

    try:
        raw = ask_ai(prompt, user=request.user, use_rag=False, model=selected_model)
        import re
        clean_raw = re.sub(r"```json\\s*|\\s*```", "", raw).strip()
        json_start = clean_raw.find("[")
        json_end = clean_raw.rfind("]") + 1
        if json_start != -1 and json_end != 0:
            clean_raw = clean_raw[json_start:json_end]
            
        questions = json.loads(clean_raw)
        return JsonResponse({"success": True, "questions": questions})
    except Exception as e:
        return JsonResponse({"error": f"NEXA Quiz Generation Failed: {str(e)}"}, status=500)

def quiz_report_ajax(request, pk):
    """AJAX: process post-quiz analytics and update StudentConceptState."""
    if request.method != "POST":
        return JsonResponse({"error": "POST required"}, status=405)
    material = get_object_or_404(StudyMaterial, pk=pk, owner=request.user)
    
    try:
        data = json.loads(request.body)
    except:
        return JsonResponse({"error": "Invalid JSON"}, status=400)
    
    results = data.get("results", [])
    # results: [{concept, isRight, confidence, error_type, chosen}]
    
    concept_stats = {}
    total_calibration_correct = 0
    total_calibration_count = 0
    
    for r in results:
        concept_name = r.get("concept", "General")
        is_right = r.get("isRight", False)
        confidence = r.get("confidence", 3)  # 1-5 scale
        error_type = r.get("error_type", "unknown")
        
        if concept_name not in concept_stats:
            concept_stats[concept_name] = {"correct": 0, "total": 0, "errors": [], "high_conf_wrong": 0}
        
        concept_stats[concept_name]["total"] += 1
        if is_right:
            concept_stats[concept_name]["correct"] += 1
        else:
            concept_stats[concept_name]["errors"].append(error_type)
        
        # Confidence calibration
        total_calibration_count += 1
        if (is_right and confidence >= 4) or (not is_right and confidence <= 2):
            total_calibration_correct += 1
        if not is_right and confidence >= 4:
            concept_stats[concept_name]["high_conf_wrong"] += 1
    
    # Calculate calibration score
    calibration_score = round((total_calibration_correct / max(total_calibration_count, 1)) * 100)
    
    # Update StudentConceptState for each concept
    from .models import ConceptNode, StudentConceptState
    from django.utils import timezone
    
    for concept_name, stats in concept_stats.items():
        try:
            concept_node = ConceptNode.objects.get(material=material, name__iexact=concept_name)
            state, created = StudentConceptState.objects.get_or_create(
                user=request.user, concept=concept_node
            )
            # Update strength (rolling average)
            new_score = round((stats["correct"] / max(stats["total"], 1)) * 100)
            if created:
                state.concept_strength = new_score
            else:
                state.concept_strength = round((state.concept_strength * 0.6) + (new_score * 0.4))
            
            # Update error profile
            existing_errors = state.error_profile or {}
            for err in stats["errors"]:
                existing_errors[err] = existing_errors.get(err, 0) + 1
            state.error_profile = existing_errors
            
            # Calibration
            state.confidence_calibration = calibration_score
            state.last_reviewed = timezone.now()
            state.save()
        except ConceptNode.DoesNotExist:
            pass  # Custom concepts without nodes are skipped
    
    # Build response with analytics
    error_breakdown = {}
    for stats in concept_stats.values():
        for err in stats["errors"]:
            error_breakdown[err] = error_breakdown.get(err, 0) + 1
    
    return JsonResponse({
        "success": True,
        "calibration_score": calibration_score,
        "concept_mastery": {k: {"accuracy": round((v["correct"]/max(v["total"],1))*100), "errors": v["errors"], "overconfident": v["high_conf_wrong"]} for k, v in concept_stats.items()},
        "error_breakdown": error_breakdown
    })


def flashcards_view(request, pk):
    """Dedicated page: AI-generated flashcards from a study material."""
    material = get_object_or_404(StudyMaterial, pk=pk, owner=request.user)
    return render(request, 'materials/flashcards.html', {
        'material': material, 'title': f'Flashcards: {material.title}'
    })


@login_required
def flashcards_ajax(request, pk):
    """AJAX: generate flashcards."""
    if request.method != 'POST':
        return JsonResponse({'error': 'POST required'}, status=405)
    material = get_object_or_404(StudyMaterial, pk=pk, owner=request.user)
    if not material.extracted_text:
        return JsonResponse({'error': 'No text could be extracted from this material.'}, status=400)
    prompt = f"""Create 15 flashcards from the study material titled "{material.title}".

Format each flashcard EXACTLY like this, separated by a blank line:
FRONT: [term or concept â€” keep it short]
BACK: [clear definition or explanation]

FRONT: [next term]
BACK: [definition]

Content:
{material.extracted_text[:4000]}

Flashcards:"""
    try:
        data = {}
        if request.body:
            try: data = json.loads(request.body)
            except: pass
        selected_model = resolve_model(data.get('model'))
        raw = ask_ai(prompt, user=request.user, use_rag=False, model=selected_model)
        import re
        pairs = re.findall(r'FRONT:\s*(.+?)\nBACK:\s*(.+?)(?=\nFRONT:|\Z)', raw, re.DOTALL)
        cards = [{'front': f.strip(), 'back': b.strip()} for f, b in pairs]
        if not cards:
            blocks = [b.strip() for b in raw.split('\n\n') if b.strip()]
            for block in blocks:
                lines = [l.strip() for l in block.split('\n') if l.strip()]
                if len(lines) >= 2:
                    cards.append({'front': lines[0], 'back': ' '.join(lines[1:])})
        return JsonResponse({'success': True, 'cards': cards})
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


@login_required
def flashcards_save(request, pk):
    """Save a generated flashcard deck."""
    if request.method != 'POST':
        return JsonResponse({'error': 'POST required'}, status=405)
    material = get_object_or_404(StudyMaterial, pk=pk, owner=request.user)
    try:
        data = json.loads(request.body)
        cards = data.get('cards', [])
        name = data.get('name', '').strip() or f"Deck {material.flashcard_decks.count() + 1}"
        if not cards:
            return JsonResponse({'error': 'No cards to save.'}, status=400)
        from .models import SavedFlashcardDeck
        deck = SavedFlashcardDeck.objects.create(
            material=material, owner=request.user, name=name, cards=cards
        )
        return JsonResponse({'success': True, 'deck_id': deck.pk, 'name': deck.name, 'created_at': deck.created_at.strftime('%b %d, %Y')})
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


@login_required
def flashcards_decks(request, pk):
    """List saved decks for a material."""
    material = get_object_or_404(StudyMaterial, pk=pk, owner=request.user)
    from .models import SavedFlashcardDeck
    decks = SavedFlashcardDeck.objects.filter(material=material, owner=request.user).values(
        'id', 'name', 'created_at'
    )
    result = [{'id': d['id'], 'name': d['name'], 'created_at': d['created_at'].strftime('%b %d, %Y'), 'count': SavedFlashcardDeck.objects.get(pk=d['id']).cards.__len__()} for d in decks]
    return JsonResponse({'decks': result})


@login_required
def flashcards_load(request, pk, deck_id):
    """Load a saved deck's cards."""
    material = get_object_or_404(StudyMaterial, pk=pk, owner=request.user)
    from .models import SavedFlashcardDeck
    deck = get_object_or_404(SavedFlashcardDeck, pk=deck_id, material=material, owner=request.user)
    return JsonResponse({'success': True, 'cards': deck.cards, 'name': deck.name})


@login_required
def flashcards_delete_deck(request, pk, deck_id):
    """Delete a saved deck."""
    if request.method != 'POST':
        return JsonResponse({'error': 'POST required'}, status=405)
    material = get_object_or_404(StudyMaterial, pk=pk, owner=request.user)
    from .models import SavedFlashcardDeck
    deck = get_object_or_404(SavedFlashcardDeck, pk=deck_id, material=material, owner=request.user)
    deck.delete()
    return JsonResponse({'success': True})


@login_required
def select_material_for_action(request, action):
    """Page to select a material before performing an action (summarize/quiz/flashcards/podcast)."""
    materials = StudyMaterial.objects.filter(owner=request.user)
    return render(request, 'materials/select_material.html', {
        'materials': materials, 'action': action, 'title': f'Select Material â€” {action.title()}'
    })


@login_required
def wiki_image_ajax(request):
    """Search Wikipedia for an image relevant to a query keyword. No API key needed."""
    if request.method != 'POST':
        return JsonResponse({'error': 'POST required'}, status=405)
    try:
        import urllib.request
        import urllib.parse
        data = json.loads(request.body)
        query = data.get('query', '').strip()
        if not query:
            return JsonResponse({'image_url': None})

        # Step 1: search Wikipedia for the best matching article
        search_url = (
            'https://en.wikipedia.org/w/api.php?action=query&list=search'
            '&srsearch=' + urllib.parse.quote(query) +
            '&srlimit=1&format=json&origin=*'
        )
        req = urllib.request.Request(search_url, headers={'User-Agent': 'NexaApp/1.0'})
        with urllib.request.urlopen(req, timeout=5) as r:
            search_data = json.loads(r.read())

        results = search_data.get('query', {}).get('search', [])
        if not results:
            return JsonResponse({'image_url': None})

        title = results[0]['title']

        # Step 2: get the page's main image
        image_url_api = (
            'https://en.wikipedia.org/w/api.php?action=query&titles='
            + urllib.parse.quote(title) +
            '&prop=pageimages&pithumbsize=600&format=json&origin=*'
        )
        req2 = urllib.request.Request(image_url_api, headers={'User-Agent': 'NexaApp/1.0'})
        with urllib.request.urlopen(req2, timeout=5) as r2:
            img_data = json.loads(r2.read())

        pages = img_data.get('query', {}).get('pages', {})
        image_url = None
        for page in pages.values():
            thumb = page.get('thumbnail', {})
            if thumb.get('source'):
                image_url = thumb['source']
                break

        return JsonResponse({'image_url': image_url, 'title': title})

    except Exception as e:
        # Silently fail â€” image is optional
        return JsonResponse({'image_url': None})



@login_required
def learn_view(request, pk):
    """The High-Performance Study Studio. Renders a shell immediately, deferring slide extraction."""
    material = get_object_or_404(StudyMaterial, pk=pk, owner=request.user)
    ext = os.path.splitext(material.file.name)[1].lower() if material.file else ''
    
    # We pass placeholders initially to avoid HUGE HTML blobs. 
    # Frontend will fetch actual data via AJAX.
    return render(request, 'materials/viewer.html', {
        'material': material,
        'pages_json': '[]', # Defer loading
        'total_pages': 0, # Update via AJAX
        'title': f'Study Studio: {material.title}',
        'file_ext': ext,
        'defer_load': True
    })


@login_required
def material_slides_api(request, pk):
    """JSON API with caching: returns processed slides for the Study Studio."""
    material = get_object_or_404(StudyMaterial, pk=pk, owner=request.user)
    
    # Return from cache if we have it!
    if material.extracted_pages_json:
        return JsonResponse({
            'success': True,
            'pages': material.extracted_pages_json,
            'total': len(material.extracted_pages_json),
            'cached': True
        })

    # First time? Extract
    ext = os.path.splitext(material.file.name)[1].lower() if material.file else ''
    pages = []
    
    file_path = None
    try:
        file_path = material.file.path
        if not os.path.exists(file_path): file_path = None
    except Exception: file_path = None

    if file_path and ext == '.pptx':
        pages = _extract_pptx_pages(file_path)
    elif file_path and ext == '.pdf':
        pages = _extract_pdf_pages(file_path)
    elif file_path and ext in ('.docx', '.doc'):
        pages = _extract_docx_pages(file_path)
    else:
        # Text-only fallback for text/markdown or where extraction failed
        text = material.extracted_text or ''
        if text:
            paragraphs = [p.strip() for p in text.split('\n') if p.strip()]
            current, size = [], 0
            for para in paragraphs:
                current.append(para)
                size += len(para)
                if size >= 600:
                    pages.append({'text': '\n\n'.join(current), 'images': [], 'label': f'Section {len(pages)+1}'})
                    current, size = [], 0
            if current:
                pages.append({'text': '\n\n'.join(current), 'images': [], 'label': f'Section {len(pages)+1}'})

    if not pages:
        pages = [{'text': 'No content could be extracted from this file.', 'images': [], 'label': 'Error'}]

    # Post-process to detect chapters/modules and assign titles
    current_chapter = None
    default_chapter = "Overview"
    
    # regex for common lecture/chapter markers
    chapter_pattern = re.compile(r'^(?:Chapter|Module|Lecture|Session|Topic|Unit|Lesson|Part)\b', re.IGNORECASE)
    
    for page in pages:
        text = page.get('text', '')
        lines = text.split('\n')
        # Check first few lines for a chapter heading
        for line in lines[:5]:
            line = line.strip()
            if not line: continue
            
            if chapter_pattern.match(line):
                current_chapter = line
                break
        
        if not current_chapter:
            current_chapter = default_chapter
        
        page['chapter_title'] = current_chapter

    # Cache it for next time!
    material.extracted_pages_json = pages
    material.save(update_fields=['extracted_pages_json'])

    return JsonResponse({
        'success': True,
        'title': material.title,
        'pages': pages,
        'total': len(pages),
        'cached': False
    })

def _extract_pptx_pages(file_path):
    """Extract each slide as a page with text + images (base64). Recursive for groups."""
    import base64, io
    pages = []
    try:
        import pptx
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = pptx.Presentation(file_path)

        def _get_shape_images(shape, img_list):
            # Recursive image extraction from groups or solo pictures
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_blob = shape.image.blob
                    img_ext = shape.image.ext or 'png'
                    b64 = base64.b64encode(img_blob).decode('utf-8')
                    mime = 'image/jpeg' if img_ext.lower() in ('jpg','jpeg') else f'image/{img_ext.lower()}'
                    img_list.append(f'data:{mime};base64,{b64}')
                except Exception: pass
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for s in shape.shapes:
                    _get_shape_images(s, img_list)
            # Check for picture fills in any shape type
            try:
                if hasattr(shape, 'fill') and hasattr(shape.fill, 'type') and shape.fill.type == 6: # MSO_FILL.PICTURE
                    img_blob = shape.fill.picture_fill.image.blob
                    img_ext = shape.fill.picture_fill.image.ext or 'png'
                    b64 = base64.b64encode(img_blob).decode('utf-8')
                    mime = 'image/jpeg' if img_ext.lower() in ('jpg','jpeg') else f'image/{img_ext.lower()}'
                    img_list.append(f'data:{mime};base64,{b64}')
            except Exception: pass

        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts = []
            images = []
            for shape in slide.shapes:
                # Text
                if hasattr(shape, 'text') and shape.text.strip():
                    text_parts.append(shape.text.strip())
                # Images
                _get_shape_images(shape, images)

            pages.append({
                'text': '\n\n'.join(text_parts) or f'Slide {slide_num}',
                'images': images,
                'label': f'Slide {slide_num}',
            })
    except ImportError:
        pages = [{'text': 'python-pptx is required. Install with: pip install python-pptx', 'images': []}]
    except Exception as e:
        pages = [{'text': f'Error reading presentation: {str(e)}', 'images': []}]
    return pages


def _extract_pdf_pages(file_path):
    """Extract each PDF page as text + images (base64). Incudes a full-page snapshot for diagrams."""
    import base64
    pages = []
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        for page_num, page in enumerate(doc, 1):
            text = page.get_text('text').strip()
            images = []
            
            # --- Extract Text and Individual Diagrams Only ---
            # (Redundant whole-page visual snapshot removed for a cleaner experience)

            # --- 2. Significant Diagrams only (skip small fragments/icons) ---
            for img_index, img in enumerate(page.get_images(full=True)):
                try:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    if base_image:
                        # Professional filter: skip very small fragments (signatures, decorative boxes)
                        width = base_image.get('width', 0)
                        height = base_image.get('height', 0)
                        if width < 120 or height < 120: continue
                        
                        img_bytes = base_image['image']
                        img_ext = base_image.get('ext', 'png')
                        b64 = base64.b64encode(img_bytes).decode('utf-8')
                        mime = 'image/jpeg' if img_ext.lower() in ('jpg','jpeg') else f'image/{img_ext.lower()}'
                        images.append(f'data:{mime};base64,{b64}')
                except Exception: pass
                
            pages.append({
                'text': text or f'Page {page_num}',
                'images': images,
                'label': f'Page {page_num}',
            })
        doc.close()
        return pages
    except ImportError:
        pass
    except Exception:
        pass

    # Fallback: PyPDF2 (text only)
    try:
        import PyPDF2
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page_num, page in enumerate(reader.pages, 1):
                text = page.extract_text() or ''
                pages.append({'text': text.strip() or f'Page {page_num}', 'images': [], 'label': f'Page {page_num}'})
    except Exception as e:
        pages = [{'text': f'Error reading PDF: {str(e)}', 'images': []}]
    return pages


def _extract_docx_pages(file_path):
    """Extract Word doc — upgraded to harvest ALL embedded images."""
    import base64, io
    pages = []
    try:
        import docx
        from docx.document import Document as _Document
        from docx.oxml.table import CT_Tbl
        from docx.oxml.text.paragraph import CT_P
        from docx.table import _Cell, Table
        from docx.text.paragraph import Paragraph

        doc = docx.Document(file_path)
        
        # Pre-harvest ALL images from the relationship parts to be safe
        fallback_images = []
        try:
            for rel in doc.part.rels.values():
                if 'image' in rel.reltype:
                    img_part = rel.target_part
                    img_blob = img_part.blob
                    ct = img_part.content_type or 'image/png'
                    b64 = base64.b64encode(img_blob).decode('utf-8')
                    fallback_images.append(f'data:{ct};base64,{b64}')
        except Exception: pass

        current_text, current_images, size = [], [], 0
        page_num = 1

        def add_page():
            nonlocal page_num, current_text, current_images, size
            if current_text or current_images:
                # If no images found for this chunk, give it access to one fallback image if available?
                # No, just include everything found in this element.
                pages.append({
                    'text': '\n\n'.join(current_text),
                    'images': current_images,
                    'label': f'Section {page_num}'
                })
                current_text, current_images, size = [], [], 0
                page_num += 1

        # Iterate through all elements in order
        for block in doc.element.body:
            if isinstance(block, CT_P):
                p = Paragraph(block, doc)
                txt = p.text.strip()
                if txt:
                    current_text.append(txt)
                    size += len(txt)
                
                # Check for images in this paragraph
                for rel in p.part.rels.values():
                    if 'image' in rel.reltype:
                        try:
                            # If it's referenced in the paragraph's XML, it belongs here
                            if rel.rId in p._element.xml:
                                img_part = rel.target_part
                                img_bytes = img_part.blob
                                ct = img_part.content_type or 'image/png'
                                b64 = base64.b64encode(img_bytes).decode('utf-8')
                                current_images.append(f'data:{ct};base64,{b64}')
                        except Exception: pass

                if size >= 800:
                    add_page()
            
            elif isinstance(block, CT_Tbl):
                t = Table(block, doc)
                for row in t.rows:
                    for cell in row.cells:
                        txt = cell.text.strip()
                        if txt:
                            current_text.append(txt)
                            size += len(txt)
                if size >= 800:
                    add_page()

        add_page() # Final chunk
        
        # Final pass: If NO images were found in the sections, but they exist in doc, add to the first section
        if not any(p['images'] for p in pages) and fallback_images:
            if pages:
                pages[0]['images'] = fallback_images

    except ImportError:
        pages = [{'text': 'python-docx is required. Install with: pip install python-docx', 'images': []}]
    except Exception as e:
        pages = [{'text': f'Error reading document: {str(e)}', 'images': []}]
    return pages


@login_required
def learn_ajax(request, pk):
    import json, re
    """AJAX: AI assist for a specific page/section of the material."""
    if request.method != 'POST':
        return JsonResponse({'error': 'POST required'}, status=405)
    material = get_object_or_404(StudyMaterial, pk=pk, owner=request.user)
    try:
        data = json.loads(request.body)
        action = data.get('action', 'summarise')
        page_text = data.get('page_text', '').strip()
        page_label = data.get('page_label', 'this section')
        selected_model = resolve_model(data.get('model'))
        thread_id = data.get('thread_id')

        # Resolve Memory (Thread & History)
        thread = None
        history = []
        
        if thread_id:
            thread = ChatThread.objects.filter(id=thread_id, user=request.user).first()
        
        # Fetch history if we have a thread
        if thread:
            recent = Conversation.objects.filter(thread=thread).order_by('-created_at')[:15]
            history = [{'message': c.message, 'response': c.response} for c in reversed(recent)]

        # [AI NORM] Unlimited context support (up to 250k chars)
        MAX_INPUT = 250000
        if len(page_text) > MAX_INPUT:
            page_text = page_text[:MAX_INPUT] + "... (Document too large for single scan)"

        if not page_text and action not in ['podcast_list', 'podcast_detail']:
            return JsonResponse({'error': 'No text provided.'}, status=400)

        elif action == 'chat':
            user_msg = data.get('message', '').strip()
            images = data.get('images', []) # Pass current slide images for visual reasoning
            
            prompt = f'''You are AI Pal, a friendly but genius tutor. A student is looking at this section titled "{page_label}".
            
            Section Content:
            {page_text}
            
            User Question: "{user_msg}"
            
            Guidelines:
            1. Use the provided section text as your primary source of truth.
            2. If images are provided in the context, you can reference them to explain diagrams or visual concepts.
            3. Be encouraging, concise, and smart. 
            4. If the student asks something unrelated to the material, gently bring them back to the topic.
            '''
            
            # If we have images, we use chat_with_image logic (needs gpt-4o-mini)
            if images:
                # We reuse the existing chat_with_image logic if possible or just call OpenAI directly here
                from .openai_utils import chat_with_image
                # Combine all images into the request for context
                context_image = images[0] # Usually the most important snapshot
                raw = chat_with_image(user_msg, context_image, system_prompt=prompt)
            else:
                raw = ask_ai(prompt, user=request.user, use_rag=False, history=history, model=selected_model)

            # Persist to Thread
            if not thread:
                thread = ChatThread.objects.create(user=request.user, title=f"Study: {material.title[:50]}")
                
            Conversation.objects.create(
                user=request.user,
                thread=thread,
                message=user_msg,
                response=raw
            )
                
            return JsonResponse({'success': True, 'action': action, 'result': raw, 'thread_id': thread.id})

        elif action == 'summarise' or action == 'brief' or action == 'explain':
            prompt = f'''Perform a Deep Intelligence Scan on this study section titled "{page_label}".
Return ONLY a valid JSON object with these keys:
{{"summary": "clear bullet points of key details", "vocab": ["term1", "term2", "term3", "term4"], "recap": "One punchy genius sentence summarizing the core takeaway", "actions": ["Do X to master this", "Compare this to Y", "Draw a diagram of Z"]}}

Section Content:
{page_text[:200000]}''' # Safe limit for GPT-4o-mini
            
            try:
                raw = ask_ai(prompt, user=request.user, use_rag=False)
                match = re.search(r'\{.*\}', raw, re.DOTALL)
                if match:
                    data = json.loads(match.group())
                    return JsonResponse({'success': True, 'action': action, 'result': data})
                
                # Fallback for plain text if AI fails JSON
                return JsonResponse({'success': True, 'action': action, 'result': {'summary': raw, 'recap': 'Summary complete.', 'vocab': [], 'actions': []}})
            except Exception as e:
                return JsonResponse({'success': False, 'error': f"NEXA Briefing Error: {str(e)}"}, status=400)

        elif action == 'quiz':
            prompt = f'''Create 4 multiple choice questions from this section titled "{page_label}".
Return ONLY valid JSON array, no extra text:
[
  {{"q":"question text","opts":{{"A":"option","B":"option","C":"option","D":"option"}},"ans":"A","explanation":"why A is correct"}},
  ...
]

Section:
{page_text}'''
            try:
                raw = ask_ai(prompt, user=request.user, use_rag=False)
                # extract JSON array
                match = re.search(r'\[.*\]', raw, re.DOTALL)
                if match:
                    questions = json.loads(match.group())
                    return JsonResponse({'success': True, 'action': action, 'questions': questions})
                return JsonResponse({'success': False, 'error': 'AI returned invalid quiz format. Try again.'}, status=500)
            except Exception as e:
                return JsonResponse({'success': False, 'error': f"NEXA Studio Quiz Error: {str(e)}"}, status=400)

        elif action == 'flashcards':
            prompt = f'''Create 6 flashcards from this section titled "{page_label}".
Return ONLY valid JSON array, no extra text:
[
  {{"front":"term or concept","back":"clear definition or explanation"}},
  ...
]

Section:
{page_text}'''
            try:
                raw = ask_ai(prompt, user=request.user, use_rag=False)
                import re
                match = re.search(r'\[.*\]', raw, re.DOTALL)
                if match:
                    cards = json.loads(match.group())
                    return JsonResponse({'success': True, 'action': action, 'cards': cards})
                # fallback: parse FRONT/BACK format
                pairs = re.findall(r'FRONT:\s*(.+?)\nBACK:\s*(.+?)(?=\nFRONT:|\Z)', raw, re.DOTALL)
                if pairs:
                    cards = [{'front': f.strip(), 'back': b.strip()} for f, b in pairs]
                    return JsonResponse({'success': True, 'action': action, 'cards': cards})
                return JsonResponse({'success': False, 'error': 'AI returned invalid card format. Try again.'}, status=500)
            except Exception as e:
                return JsonResponse({'success': False, 'error': f"NEXA Studio Cards Error: {str(e)}"}, status=400)

        elif action == 'podcast':
            # Use local function instead of incorrect import
            from .audio_utils import generate_podcast_segments
            prompt = build_podcast_prompt(page_text, page_label)
            try:
                script_text = ask_ai(prompt, user=request.user, use_rag=False)
                # Generate natural audio segments (Kojo/Afia)
                segments = generate_podcast_segments(script_text, pk)
                return JsonResponse({
                    'success': True, 
                    'action': action, 
                    'script_text': script_text,
                    'script_json': segments
                })
            except Exception as e:
                return JsonResponse({'success': False, 'error': f"Podcast Error: {str(e)}"}, status=500)

        elif action == 'podcast_call':
            from .audio_utils import generate_podcast_segments
            question = data.get('question', '').strip()
            if not question:
                return JsonResponse({'error': 'No question provided.'}, status=400)
                
            prompt = f'''You are directing "The Alex & Sam Masterclass".
The podcast was just interrupted because a listener called in with a question!
The current topic is: {page_label}

Listener's Question: "{question}"

Write a short, pure dialogue script addressing the question.
STRICT RULES:
1. START EXACTLY LIKE THIS:
   Sam: We actually have a caller on the line! Caller, you there?
   Alex: That is a fascinating question.
2. HAVE THEM ANSWER IT NATURALLY.
3. END WITH:
   Sam: Thanks for calling in! Now, back to what we were saying...
4. NO BRACKETS, NO REPETITIVE NAMES.

Now answer:'''
            try:
                script_text = ask_ai(prompt, user=request.user, use_rag=False)
                segments = generate_podcast_segments(script_text, pk)
                return JsonResponse({
                    'success': True, 
                    'action': action, 
                    'script_json': segments
                })
            except Exception as e:
                return JsonResponse({'success': False, 'error': f"Podcast Call Error: {str(e)}"}, status=500)

        elif action == 'podcast_save':
            from .audio_utils import generate_podcast_segments
            script_json = data.get('script_json')
            script_text = data.get('script_text', '').strip()
            name = data.get('name', f"Podcast for {material.title}").strip()
            
            # If we already have the segments (from a Live run), save them directly
            if script_json:
                segments = script_json
            elif script_text:
                segments = generate_podcast_segments(script_text, material.pk)
            else:
                return JsonResponse({'success': False, 'error': 'No content to save.'})
            
            pod = SavedPodcast.objects.create(
                material=material,
                owner=request.user,
                name=name,
                script_json=segments
            )
            return JsonResponse({'success': True, 'podcast_id': pod.id})

        elif action == 'podcast_list':
            pods = SavedPodcast.objects.filter(material=material).values('id', 'name', 'created_at')
            pod_list = []
            for p in pods:
                pod_list.append({
                    'id': p['id'],
                    'name': p['name'],
                    'created_at': p['created_at'].strftime('%Y-%m-%d %H:%M')
                })
            return JsonResponse({'success': True, 'podcasts': pod_list})

        elif action == 'podcast_detail':
            pod_id = data.get('podcast_id')
            pod = get_object_or_404(SavedPodcast, id=pod_id, owner=request.user)
            return JsonResponse({'success': True, 'script_json': pod.script_json})

        elif action == 'chat' or action == 'help' or action == 'explain' or action == 'qa':
            # Standard AI assist (non-multi-modal) with history support
            user_msg = data.get('message', '').strip() or f"Help me with {page_label}"
            prompt = f'''You are Kojo, a brilliant and supportive study assistant. 
I am studying this section: "{page_label}".

CONTENT:
{page_text}

STUDENT REQUEST:
{user_msg}'''
            
            raw = ask_ai(prompt, user=request.user, use_rag=False, history=history, model=selected_model)
            
            # Save to Thread
            if not thread:
                thread = ChatThread.objects.create(user=request.user, title=f"Study: {material.title[:50]}")
            
            Conversation.objects.create(
                user=request.user,
                thread=thread,
                message=user_msg,
                response=raw
            )
            
            return JsonResponse({'success': True, 'action': action, 'result': raw, 'thread_id': thread.id})

        elif action == 'podcast_question':
            from .audio_utils import generate_podcast_segments
            user_q = data.get('question', '').strip()
            conv_context = data.get('conversation_context', '').strip()
            prompt = f'''You are Sam and Alex, hosts of a study podcast. A student just "called in" with a question about this section: "{page_label}".

Podcast Conversation So Far:
{conv_context}

Caller's Question: "{user_q}"

Respond as Sam or Alex in a friendly, conversational way. Answer the question accurately using the section text provided AND referencing the discussion above if relevant. 
After answering, say exactly: "Alright, back to the show!" to signal a transition.

Section Content:
{page_text}'''
            raw = ask_ai(prompt, user=request.user, use_rag=False)
            
            # Generate premium audio for the answer
            segments = generate_podcast_segments(raw, material.pk)
            audio_url = segments[0]['audio_url'] if segments else None
            
            return JsonResponse({'success': True, 'action': action, 'result': raw, 'audio_url': audio_url})

        elif action == 'song_generate':
            from .audio_utils import generate_podcast_segments
            genre = data.get('genre', 'Lofi').strip()
            prompt = f'''You are a world-class Divine Singer. Create catchy, rhythmic, and rhymed song lyrics about this study section titled "{page_label}" in the style of {genre} music.

Lyrics Rules:
1. Every line MUST rhyme with the next (Strong AABB or ABAB schemes).
2. The rhythm should be musical and easy to "sing" or "rap" over a beat.
3. EVERY LINE of the lyrics must be prefixed with "Divine: ".
4. Focus on making the chorus very catchy and repetitive.
5. Use section headers like [Intro], [Verse 1], [Chorus], [Verse 2], [Chorus], [Outro] on their own lines.

Example:
[Verse 1]
Divine: We're diving deep into the cell today
Divine: Working on our notes in a brand new way!

Section Content:
{page_text}'''
            raw_lyrics = ask_ai(prompt, user=request.user, use_rag=False)
            
            # Bake the "Singing" vocals using our premium neural engine
            segments = generate_podcast_segments(raw_lyrics, material.pk, use_premium=True)
            return JsonResponse({'success': True, 'action': action, 'lyrics_json': segments, 'raw_text': raw_lyrics})

        elif action == 'song_save':
            lyrics_json = data.get('lyrics_json')
            genre = data.get('genre', 'Lofi').strip()
            name = data.get('name', f"{genre} Beat about {material.title}").strip()
            
            song = SavedStudySong.objects.create(
                material=material,
                owner=request.user,
                name=name,
                genre=genre,
                lyrics_json=lyrics_json
            )
            return JsonResponse({'success': True, 'song_id': song.id})

        elif action == 'song_list':
            songs = SavedStudySong.objects.filter(material=material, owner=request.user)
            song_data = [{'id': s.id, 'name': s.name, 'genre': s.genre, 'created_at': s.created_at.strftime('%b %d')} for s in songs]
            return JsonResponse({'success': True, 'songs': song_data})

        elif action == 'song_detail':
            song_id = data.get('song_id')
            song = get_object_or_404(SavedStudySong, id=song_id, owner=request.user)
            return JsonResponse({'success': True, 'lyrics_json': song.lyrics_json, 'genre': song.genre})

        elif action == 'podcast':
            from .audio_utils import generate_podcast_segments
            prompt = f'''Create an immersive, high-energy podcast script explaining this section titled "{page_label}".
            
Hosts: Alex (energetic, funny) and Sam (brilliant, grounded).
Speech Cues:
- Use [Excited] for big realizations.
- Use [Whisper] or [Serious] for key secrets or deep points.
- Use [Slower] for complex explanations.
- Use interjections like "Oh wow!", "Wait, really?", "Hmm...", and "Aha!"
- Keep it under 3 minutes. Plain text only. Use Alex: and Sam: roles.

Section:
{page_text}'''
            try:
                raw = ask_ai(prompt, user=request.user, use_rag=False)
                # Auto-generate high-quality audio segments for Live playback!
                segments = generate_podcast_segments(raw, material.pk)
                return JsonResponse({'success': True, 'action': action, 'script_json': segments})
            except Exception as e:
                # If the AI call fails (Context Window or other), provide a fallback message
                err_msg = str(e)
                if "context length" in err_msg.lower():
                    err_msg = "This chapter is too massive for a single podcast! Try breaking it into smaller sections."
                return JsonResponse({'success': False, 'error': f"NEXA Studio: {err_msg}"}, status=400)

        elif action == 'marginalia':
            prompt = f'''Generate 2-3 short, clever, and helpful study hints or "scribbles" for this section: "{page_label}".
Each hint should be max 10 words, informal student-speak. Use emojis.
Format:
HINT: [text]
HINT: [text]

Section Content:
{page_text}'''
            try:
                raw = ask_ai(prompt, user=request.user, use_rag=False)
                hints = [h.replace('HINT:', '').strip() for h in raw.split('\n') if 'HINT:' in h]
                return JsonResponse({'success': True, 'action': action, 'hints': hints[:3]})
            except:
                return JsonResponse({'success': False, 'error': 'Failed to generate marginalia.'}, status=500)

        elif action == 'paste_page':
            # Append a custom page to the material's notebook
            content = data.get('content', '').strip()
            label = data.get('label', 'Saved Note').strip()
            if not content: return JsonResponse({'success': False})
            
            pages = material.extracted_pages_json or []
            new_page = {
                'text': content,
                'images': [],
                'label': label,
                'chapter_title': f"Saved: {label}"
            }
            pages.append(new_page)
            material.extracted_pages_json = pages
            material.save(update_fields=['extracted_pages_json'])
            return JsonResponse({'success': True, 'page_index': len(pages)-1})

        elif action == 'global_briefing':
            # Summarize the 5 most recent study materials
            from .audio_utils import generate_podcast_segments
            recent_materials = StudyMaterial.objects.filter(owner=request.user).order_by('-created_at')[:5]
            if not recent_materials:
                return JsonResponse({'success': False, 'error': 'No study materials found to brief.'})
            
            titles = ", ".join([m.title for m in recent_materials])
            prompt = f'''You are Sam and Alex from NEXA Intelligence. Create a "Morning Briefing" podcast script (max 2 mins).
Briefly mention the 5 latest materials the student has been working on: {titles}.

Speech Style:
- Use [Excited] when mentioning progress.
- Use [Thoughtful] for study tips.
- Use [Fast] for quick recaps.
- Be extremely motivating! Provide one "Genius Tip" for each or a general overview. 
Plain text only. Use Alex: and Sam: roles.'''
            
            try:
                raw = ask_ai(prompt, user=request.user, use_rag=False)
                segments = generate_podcast_segments(raw, 0) # Global ID
                return JsonResponse({'success': True, 'action': action, 'script_json': segments, 'raw_text': raw})
            except Exception as e:
                return JsonResponse({'success': False, 'error': str(e)}, status=500)

        elif action == 'full_podcast':
            # Podcast for the ENTIRE document (using extracted_text)
            from .audio_utils import generate_podcast_segments
            text = material.extracted_text or ""
            if not text: return JsonResponse({'success': False, 'error': 'No text found in document.'})
            
            # No-limit expansion: Supporting massive full docs
            if len(text) > 250000: text = text[:250000] + "..."
            
            prompt = f'''Create an ultimate, comprehensive study podcast for the FULL document: "{material.title}".
Alex and Sam should dive deep into the major themes and key takeaways. Make it feel like a professional masterclass.
Max 5 minutes of dialogue. Plain text only. Use Alex: and Sam: roles.

Content:
{text}'''
            try:
                raw = ask_ai(prompt, user=request.user, use_rag=False)
                segments = generate_podcast_segments(raw, material.pk)
                return JsonResponse({'success': True, 'action': action, 'script_json': segments})
            except Exception as e:
                return JsonResponse({'success': False, 'error': str(e)}, status=500)

        else:
            prompt = f'Provide general help for the section "{page_label}":\n\n{page_text}'
            raw = ask_ai(prompt, user=request.user, use_rag=False)
            return JsonResponse({'success': True, 'action': action, 'result': raw})
            return JsonResponse({'success': True, 'action': action, 'result': raw})

    except Exception as e:
        import logging
        logging.getLogger(__name__).error(f"SYSTEM_CRASH_LEARN_AJAX: {e}")


@login_required
def generate_concept_graph_ajax(request, pk):
    """AJAX endpoint to build the Concept Dependency Graph for a material."""
    if request.method != 'POST':
        return JsonResponse({'error': 'POST required'}, status=405)
        
    material = get_object_or_404(StudyMaterial, pk=pk, owner=request.user)
    
    # If concepts already exist, return them
    if material.concepts.exists():
        concepts_data = []
        for c in material.concepts.all():
            concepts_data.append({
                'id': c.id,
                'name': c.name,
                'definition': c.definition,
                'prerequisites': list(c.prerequisites.values_list('id', flat=True))
            })
        return JsonResponse({'success': True, 'concepts': concepts_data})
        
    if not material.extracted_text:
        return JsonResponse({'error': 'No text available for concept extraction.'}, status=400)

    prompt = f"""You are a Cognitive Extraction Engine. Your task is to analyze the following study material and extract the core concepts as a Dependency Graph.
    
    Analyze this text: {material.extracted_text[:6000]}
    
    Return a strictly valid JSON object (NO Markdown, NO backticks) mapping each concept.
    Format exactly like this example:
    {{
      "concepts": [
        {{
          "name": "Photosynthesis",
          "definition": "The process by which green plants and some other organisms use sunlight to synthesize nutrients from carbon dioxide and water.",
          "prerequisites": ["Chloroplasts", "Light Energy"]
        }}
      ]
    }}
    
    Extract 5 to 10 key concepts and identify their prerequisites (if any). Ensure prerequisite names exactly match other extracted concept names.
    JSON ONLY:"""
    
    try:
        raw_response = ask_ai(prompt, user=request.user, use_rag=False)
        # Strip markdown if AI accidentally includes it
        import re
        raw_response = re.sub(r'```json\s*|\s*```', '', raw_response).strip()
        
        # Try to find JSON object directly if trailing text exists
        json_start = raw_response.find('{')
        json_end = raw_response.rfind('}') + 1
        if json_start != -1 and json_end != 0:
            raw_response = raw_response[json_start:json_end]
            
        data = json.loads(raw_response)
        
        # Phase 1: Create all nodes
        concept_dict = {}
        for item in data.get('concepts', []):
            name = item.get('name', '').strip()
            if name:
                node = ConceptNode.objects.create(
                    material=material,
                    name=name[:255],
                    definition=item.get('definition', '')
                )
                concept_dict[name.lower()] = node
                
        # Phase 2: Link prerequisites
        for item in data.get('concepts', []):
            name = item.get('name', '').strip().lower()
            if name in concept_dict:
                node = concept_dict[name]
                for prereq in item.get('prerequisites', []):
                    prereq_lower = prereq.strip().lower()
                    if prereq_lower in concept_dict and prereq_lower != name:
                        node.prerequisites.add(concept_dict[prereq_lower])
        
        # Serialize response
        concepts_data = []
        for c in material.concepts.all():
            concepts_data.append({
                'id': c.id,
                'name': c.name,
                'definition': c.definition,
                'prerequisites': list(c.prerequisites.values_list('id', flat=True))
            })
        return JsonResponse({'success': True, 'concepts': concepts_data})
        
    except Exception as e:
        import traceback
        return JsonResponse({'error': f"Failed to build concept graph: {str(e)}"}, status=500)


@login_required
def grade_concept_recall_ajax(request):
    """AJAX endpoint for the Illusion of Mastery Breaker. Grades student recall."""
    if request.method != 'POST':
        return JsonResponse({'error': 'POST required'}, status=405)
        
    try:
        data = json.loads(request.body)
        concept_id = data.get('concept_id')
        student_answer = data.get('answer', '')
        time_taken = data.get('time_taken', 0.0) # seconds
        self_confidence = data.get('confidence', 50) # 0-100
        
        concept = get_object_or_404(ConceptNode, pk=concept_id)
        
        # Initialize or get state
        state, _ = StudentConceptState.objects.get_or_create(
            user=request.user,
            concept=concept
        )
        
        # Grade the answer using AI
        prompt = f"""You are the Nexa Cognitive Engine evaluating a student's active recall for the Illusion of Mastery Breaker.
        
        Concept: {concept.name}
        Actual Definition/Key Elements: {concept.definition}
        
        Student's Recall Attempt: "{student_answer}"
        
        Return a strictly valid JSON response (NO markdown formatting) with the following structure:
        {{
            "score": [0-100 score based on accuracy and conceptual understanding. If completely wrong or blank, 0],
            "feedback": "[A brief, encouraging explanation of what they missed or got right]",
            "error_type": "[One of: 'Misconception', 'Careless', 'Blank/Guess', 'None']"
        }}
        """
        
        raw_response = ask_ai(prompt, user=request.user, use_rag=False)
        # Clean json
        import re
        raw_response = re.sub(r'```json\s*|\s*```', '', raw_response).strip()
        
        json_start = raw_response.find('{')
        json_end = raw_response.rfind('}') + 1
        if json_start != -1 and json_end != 0:
            raw_response = raw_response[json_start:json_end]
            
        grading_data = json.loads(raw_response)
        
        score = grading_data.get('score', 0)
        feedback = grading_data.get('feedback', '')
        error_type = grading_data.get('error_type', 'None')
        
        # Update Cognitive Profile metrics
        # 1. Concept Strength (rolling average weighted towards latest)
        if state.concept_strength == 0:
            state.concept_strength = score
        else:
            state.concept_strength = int(state.concept_strength * 0.4 + score * 0.6)
            
        # 2. Confidence Calibration (Gap between perceived and actual)
        state.confidence_calibration = abs(self_confidence - score)
        
        # 3. Response Velocity (Average time tracking)
        if state.response_velocity == 0:
            state.response_velocity = time_taken
        else:
            state.response_velocity = (state.response_velocity + time_taken) / 2
            
        # 4. Error Profile update
        from django.utils import timezone
        state.last_reviewed = timezone.now()
        
        error_dict = state.error_profile
        error_dict[error_type] = error_dict.get(error_type, 0) + 1
        state.error_profile = error_dict
        
        state.save()
        
        return JsonResponse({
            'success': True,
            'score': score,
            'feedback': feedback,
            'new_strength': state.concept_strength,
            'error_type': error_type,
            'prerequisites_failed': list(concept.prerequisites.values_list('name', flat=True)) if score < 60 else []
        })
        
    except Exception as e:
        import traceback
        return JsonResponse({'error': str(e)}, status=500)
